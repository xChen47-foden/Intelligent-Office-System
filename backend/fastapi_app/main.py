import redis
from fastapi import FastAPI, HTTPException, Body, Depends, Request, UploadFile, File, Query, Path, APIRouter, WebSocket, WebSocketDisconnect, Form
from fastapi.middleware.cors import CORSMiddleware  # 新增CORS支持
from pydantic import BaseModel
import random
import smtplib
from email.mime.text import MIMEText
from email.header import Header
from email.utils import formataddr
import re
from sqlalchemy import Column, Integer, String, ForeignKey, DateTime, Text, Table, func, or_
from sqlalchemy.ext.declarative import declarative_base
from sqlalchemy.orm import sessionmaker, relationship
from sqlalchemy.future import select
from sqlalchemy import text
from sqlalchemy.ext.asyncio import AsyncSession, create_async_engine
import bcrypt
import jwt
from datetime import datetime, timedelta
import sys
import os
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
from routes import search, contact, group, message, assistant_upload
import shutil
import time
import json
import uuid
import base64
from fastapi.responses import FileResponse, JSONResponse
import mimetypes
import html
import httpx
import string
import requests
import asyncio
from fastapi.staticfiles import StaticFiles
import threading
from fastapi import status
from fastapi.responses import JSONResponse
from db import Base, User, SessionLocal, engine, TodayTask
# from auth_utils import verify_token  # 使用本地定义的 verify_token
from starlette.websockets import WebSocketDisconnect
from sqlalchemy import delete as sqlalchemy_delete

# 初始化FastAPI应用
app = FastAPI()

# 添加CORS中间件（关键修改）
app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:3006", "http://localhost:3007", "*"],  # 允许前端和本地调试
    allow_credentials=True,
    allow_methods=["*"],  # 允许所有方法
    allow_headers=["*"],  # 允许所有头
)

# Redis配置（可选，如果 Redis 不可用则使用内存存储）
try:
    r = redis.Redis(host='localhost', port=6379, db=0, decode_responses=True)
    r.ping()  # 测试连接
    print("Redis 连接成功")
except Exception as e:
    print(f"Redis 连接失败，将使用内存存储: {e}")
    # 使用内存字典作为备用存储
    class MemoryStore:
        def __init__(self):
            self.data = {}
        def get(self, key):
            return self.data.get(key)
        def set(self, key, value, ex=None):
            self.data[key] = value
        def delete(self, key):
            if key in self.data:
                del self.data[key]
    r = MemoryStore()

# 数据库配置
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
MEETING_DATABASE_URL = f"sqlite+aiosqlite:///{os.path.join(BASE_DIR, '../meetings.db')}"
meeting_engine = create_async_engine(MEETING_DATABASE_URL, echo=True, future=True)
MeetingSessionLocal = sessionmaker(bind=meeting_engine, class_=AsyncSession, expire_on_commit=False)

# JWT配置
SECRET_KEY = "your_secret_key"
ALGORITHM = "HS256"

# ========== 数据模型 ==========
# 会议数据模型
class Meeting(Base):
    __tablename__ = "meetings"
    id = Column(Integer, primary_key=True, index=True)
    title = Column(String)
    host = Column(String)
    time = Column(String)  # 建议用 DateTime 类型，前端传字符串
    location = Column(String)
    period = Column(String, default="")
    status = Column(String, default="upcoming")
    user_id = Column(Integer, index=True)  # 新增字段
    participants = Column(String, default="")  # 参会人列表，存储为逗号分隔的用户ID
    host_user_id = Column(Integer, default=0)  # 主持人用户ID

# 录音历史数据模型
class RecordingHistory(Base):
    __tablename__ = "recording_history"
    id = Column(Integer, primary_key=True, index=True)
    meeting_id = Column(Integer, index=True)  # 关联的会议ID
    meeting_title = Column(String)  # 会议标题
    user_id = Column(Integer, index=True)  # 录音者ID
    filename = Column(String)  # 录音文件名
    file_path = Column(String)  # 录音文件路径
    transcript = Column(Text)  # 语音转录内容
    minutes = Column(Text)  # 生成的会议纪要
    duration = Column(Integer, default=0)  # 录音时长（秒）
    file_size = Column(Integer, default=0)  # 文件大小（字节）
    created_at = Column(String)  # 创建时间
    status = Column(String, default="completed")  # 状态：recording, completed, failed

# 通知数据模型
# 删除通知数据模型

# ========== Pydantic请求模型 ==========
# 验证码请求模型
class CaptchaRequest(BaseModel):
    contact: str

class CaptchaVerifyRequest(BaseModel):
    contact: str
    captcha: str

class RegisterRequest(BaseModel):
    username: str
    password: str
    contact: str
    captcha: str
    department: str  # 新增部门字段

class LoginRequest(BaseModel):
    username: str
    password: str
    department: str  # 新增部门字段

class UpdateUserInfoRequest(BaseModel):
    realName: str = ""
    nickName: str = ""

class UpdatePreferencesRequest(BaseModel):
    theme: str
    language: str

class MeetingCreateRequest(BaseModel):
    title: str
    host: str
    time: str
    location: str
    period: str = ""
    status: str = "upcoming"
    participants: list[int] = []  # 参会人用户ID列表

class ResetPasswordRequest(BaseModel):
    contact: str
    captcha: str
    new_password: str

# 录音历史相关模型
class RecordingUploadRequest(BaseModel):
    meeting_id: int = None
    meeting_title: str = ""
    filename: str
    transcript: str = ""
    minutes: str = ""
    duration: int = 0
    file_size: int = 0

class RecordingUpdateRequest(BaseModel):
    transcript: str = ""
    minutes: str = ""

# 删除通知相关Pydantic模型

# ========== 工具函数 ==========
# 发送邮件验证码
def send_email_code(to_email, code):
    email_user = os.getenv("EMAIL_USER", "1642256761@qq.com")
    email_pass = os.getenv("EMAIL_PASS", "izeltputkiwldgba")
    email_host = os.getenv("EMAIL_HOST", "smtp.qq.com")
    email_port = int(os.getenv("EMAIL_PORT", 465))

    msg = MIMEText(
        f"您的验证码是：{code}，5分钟内有效。",
        'plain',
        'utf-8'
    )
    msg['Subject'] = Header('邮箱验证码', 'utf-8')
    msg['From'] = formataddr(('系统管理员', email_user))
    msg['To'] = to_email

    try:
        server = smtplib.SMTP_SSL(email_host, email_port)
        server.login(email_user, email_pass)
        server.sendmail(email_user, [to_email], msg.as_string())
        server.quit()
        return True
    except smtplib.SMTPAuthenticationError as e:
        raise Exception("邮箱认证失败，请检查邮箱和授权码是否正确")
    except smtplib.SMTPConnectError as e:
        raise Exception("无法连接到邮件服务器，请检查网络连接")
    except smtplib.SMTPException as e:
        raise Exception(f"邮件发送失败: {str(e)}")
    except Exception as e:
        raise Exception(f"发送邮件时发生未知错误: {str(e)}")

# ========== 认证相关API ==========
# 发送验证码接口
@app.post("/api/send-captcha")
async def send_captcha(data: CaptchaRequest):
    # 验证邮箱格式
    email_reg = r"^[\w\.-]+@[\w\.-]+\.[a-zA-Z]{2,}$"
    if not re.match(email_reg, data.contact):
        return {"code": 1, "msg": "请输入正确的邮箱格式"}
    
    # 生成验证码
    code = str(random.randint(100000, 999999))
    
    try:
        # 存储验证码到Redis
        r.setex(f"captcha:{data.contact}", 300, code)
        
        # 发送邮件
        send_email_code(data.contact, code)
        return {"code": 0, "msg": "验证码已发送"}
    except Exception as e:
        # 如果发送失败，删除Redis中的验证码
        try:
            r.delete(f"captcha:{data.contact}")
        except:
            pass
        return {"code": 500, "msg": str(e)}

# 验证验证码接口
@app.post("/api/verify-captcha")
async def verify_captcha(data: CaptchaVerifyRequest):
    real_code = r.get(f"captcha:{data.contact}")
    if not real_code:
        return {"code": 1, "msg": "验证码已过期或不存在"}
    if data.captcha != real_code:
        return {"code": 2, "msg": "验证码错误"}
    r.delete(f"captcha:{data.contact}")
    return {"code": 0, "msg": "验证成功"}

# 生成随机用户名
def generate_random_username():
    return "user_" + ''.join(random.choices(string.digits, k=5))

# 生成随机头像URL
def generate_random_avatar():
    """生成随机头像"""
    # 预设的头像列表（包括SVG和其他格式）
    avatar_list = [
        "avatar1.svg", "avatar2.svg", "avatar3.svg", 
        "avatar1.jpg", "avatar2.jpg", "avatar3.jpg", "avatar4.jpg", "avatar5.jpg",
        "avatar6.jpg", "avatar7.jpg", "avatar8.jpg", "avatar9.jpg", "avatar10.jpg",
        "2.jpg"  # 现有的头像文件
    ]
    
    # 如果有预设头像文件，随机选择一个
    uploads_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'uploads', 'avatar')
    os.makedirs(uploads_dir, exist_ok=True)
    
    # 检查是否有可用的头像文件
    available_avatars = []
    for avatar in avatar_list:
        avatar_path = os.path.join(uploads_dir, avatar)
        if os.path.exists(avatar_path):
            available_avatars.append(avatar)
    
    if available_avatars:
        return f"avatar/{random.choice(available_avatars)}"
    
    # 如果没有预设头像，生成基于字符的头像URL
    # 使用DiceBear或其他头像生成服务
    avatar_styles = ['adventurer', 'avataaars', 'big-ears', 'big-ears-neutral', 'big-smile', 'bottts', 'croodles', 'fun-emoji', 'icons', 'identicon', 'initials', 'lorelei', 'micah', 'miniavs', 'open-peeps', 'personas', 'pixel-art', 'shapes', 'thumbs']
    
    # 生成一个随机种子
    seed = ''.join(random.choices(string.ascii_lowercase + string.digits, k=10))
    style = random.choice(avatar_styles)
    
    # 返回DiceBear API URL
    return f"https://api.dicebear.com/7.x/{style}/svg?seed={seed}&backgroundColor=b6e3f4,c0aede,d1d4f9,ffd5dc,ffdfbf"

# 从URL下载头像并保存到本地
def save_avatar_from_url(avatar_url: str, user_id: int) -> str:
    try:
        import requests
        response = requests.get(avatar_url, timeout=10)
        if response.status_code == 200:
            uploads_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'uploads', 'avatar')
            os.makedirs(uploads_dir, exist_ok=True)
            
            # 保存为SVG文件
            filename = f"user_{user_id}_{int(time.time())}.svg"
            file_path = os.path.join(uploads_dir, filename)
            
            with open(file_path, 'wb') as f:
                f.write(response.content)
            
            return f"avatar/{filename}"
    except Exception as e:
        print(f"保存头像失败: {e}")
    
    # 如果保存失败，返回默认头像路径
    return "avatar/default.svg"

# 用户注册接口
@app.post("/api/register")
async def register_user(data: RegisterRequest):
    try:
        print("收到注册请求：", data.username, data.contact, data.department)
        
        # 验证必填字段
        if not data.contact or not data.contact.strip():
            raise HTTPException(400, "联系方式不能为空")
        if not data.password or not data.password.strip():
            raise HTTPException(400, "密码不能为空")
        if len(data.password) < 6:
            raise HTTPException(400, "密码长度至少为6位")
        if not data.captcha or not data.captcha.strip():
            raise HTTPException(400, "验证码不能为空")
        if not data.department or not data.department.strip():
            raise HTTPException(400, "请选择部门")
        
        # 验证验证码
        real_code = r.get(f"captcha:{data.contact}")
        print("验证码从redis获取：", real_code)
        if not real_code:
            print("验证码已过期或不存在")
            raise HTTPException(400, "验证码已过期或不存在，请重新获取")
        if data.captcha != real_code:
            print("验证码错误")
            raise HTTPException(400, "验证码错误，请重新输入")
        r.delete(f"captcha:{data.contact}")

        # 自动生成账号
        username = data.username.strip() if data.username else ""
        if not username:
            username = generate_random_username()
        async with SessionLocal() as session:
            # 检查用户名是否已存在
            result = await session.execute(select(User).where(User.username == username))
            if result.scalar_one_or_none():
                print("用户名已存在")
                raise HTTPException(400, "用户名已存在")
            
            # 检查联系方式是否已存在
            result = await session.execute(select(User).where(User.contact == data.contact))
            if result.scalar_one_or_none():
                print("联系方式已存在")
                raise HTTPException(400, "该邮箱已被注册，请使用其他邮箱")
            
            hashed_password = bcrypt.hashpw(data.password.encode(), bcrypt.gensalt()).decode()
            
            # 生成随机头像
            avatar_url = generate_random_avatar()
            
            new_user = User(
                username=username,
                password=hashed_password,
                contact=data.contact,
                real_name="",
                nick_name="",
                avatar=avatar_url,  # 使用生成的头像
                department=data.department
            )
            session.add(new_user)
            await session.commit()
            await session.refresh(new_user)  # 获取用户ID
            
            # 如果头像是外部URL，尝试下载并保存到本地
            if avatar_url.startswith('https://'):
                try:
                    local_avatar = save_avatar_from_url(avatar_url, new_user.id)
                    new_user.avatar = local_avatar
                    await session.commit()
                except Exception as e:
                    print(f"下载头像失败，使用默认头像: {e}")
                    new_user.avatar = "avatar/default.svg"
                    await session.commit()
            
            print("注册成功，已生成随机头像")
            return {"code": 0, "msg": "注册成功", "username": username, "department": data.department}
    except HTTPException:
        # 重新抛出 HTTPException，让 FastAPI 处理
        raise
    except Exception as e:
        print(f"注册过程中发生错误: {e}")
        import traceback
        traceback.print_exc()
        raise HTTPException(500, f"注册失败: {str(e)}")

# 用户登录接口
@app.post("/auth/login")
async def login(data: LoginRequest):
    try:
        print(f"[登录] 收到登录请求：用户名={data.username}, 部门={data.department}")
        
        if not data.username or not data.password or not data.department:
            return JSONResponse(
                status_code=400,
                content={"code": 1, "msg": "用户名、密码和部门不能为空"}
            )
        
        async with SessionLocal() as session:
            try:
                result = await session.execute(
                    select(User).where(
                        User.username == data.username, 
                        User.department == data.department
                    )
                )
                user = result.scalar_one_or_none()
            except Exception as db_error:
                print(f"[登录] 数据库查询错误: {db_error}")
                import traceback
                traceback.print_exc()
                return JSONResponse(
                    status_code=500,
                    content={"code": 500, "msg": f"数据库查询失败: {str(db_error)}"}
                )
            
            if not user:
                print(f"[登录] 用户不存在：{data.username}, 部门：{data.department}")
                return JSONResponse(
                    status_code=200,
                    content={"code": 1, "msg": "用户名、密码或部门错误"}
                )
            
            # 验证密码
            try:
                if not bcrypt.checkpw(data.password.encode(), user.password.encode()):
                    print(f"[登录] 密码错误：{data.username}")
                    return JSONResponse(
                        status_code=200,
                        content={"code": 1, "msg": "用户名、密码或部门错误"}
                    )
            except Exception as pwd_error:
                print(f"[登录] 密码验证错误: {pwd_error}")
                return JSONResponse(
                    status_code=500,
                    content={"code": 500, "msg": f"密码验证失败: {str(pwd_error)}"}
                )
            
            # 检查用户是否有头像，如果没有则生成一个
            if not user.avatar or user.avatar == "":
                try:
                    avatar_url = generate_random_avatar()
                    user.avatar = avatar_url
                    
                    # 如果头像是外部URL，尝试下载并保存到本地
                    if avatar_url.startswith('https://'):
                        try:
                            local_avatar = save_avatar_from_url(avatar_url, user.id)
                            user.avatar = local_avatar
                        except Exception as e:
                            print(f"[登录] 下载头像失败，使用默认头像: {e}")
                            user.avatar = "avatar/default.svg"
                    
                    await session.commit()
                    print(f"[登录] 为用户 {user.username} 生成了随机头像: {user.avatar}")
                except Exception as avatar_error:
                    print(f"[登录] 生成头像失败: {avatar_error}")
                    # 头像生成失败不影响登录，继续执行
            
            # 生成 JWT token
            try:
                payload = {
                    "sub": str(user.id),
                    "exp": datetime.utcnow() + timedelta(hours=2)
                }
                token = jwt.encode(payload, SECRET_KEY, ALGORITHM)
            except Exception as token_error:
                print(f"[登录] Token 生成失败: {token_error}")
                return JSONResponse(
                    status_code=500,
                    content={"code": 500, "msg": f"Token 生成失败: {str(token_error)}"}
                )
            
            print(f"[登录] 登录成功：{user.username} (ID: {user.id})")
            return {
                "code": 0,
                "msg": "登录成功",
                "data": {
                    "username": user.username,
                    "token": token,
                    "avatar": user.avatar or ""
                }
            }
    except HTTPException:
        raise
    except Exception as e:
        print(f"[登录] 登录异常：{e}")
        import traceback
        traceback.print_exc()
        return JSONResponse(
            status_code=500,
            content={"code": 500, "msg": f"登录异常: {str(e)}"}
        )

# ========== 认证中间件 ==========
# JWT Token验证函数
def verify_token(request: Request):
    auth = request.headers.get("Authorization")
    if not auth or not auth.startswith("Bearer "):
        print("未认证，未带token")
        raise HTTPException(401, "未认证")
    try:
        payload = jwt.decode(auth.split(" ")[1], SECRET_KEY, algorithms=[ALGORITHM])
        print("token解码成功，payload:", payload)
        # 兼容所有接口，返回 userId 和 sub
        return {
            "userId": int(payload["sub"]),
            "sub": payload["sub"]
        }
    except jwt.ExpiredSignatureError:
        print("令牌已过期")
        raise HTTPException(401, "令牌已过期")
    except jwt.InvalidTokenError as e:
        print("令牌无效:", e)
        raise HTTPException(401, "无效令牌")

# ========== 用户信息API ==========
# 测试受保护路由
@app.get("/api/protected")
async def protected_route(payload: dict = Depends(verify_token)):
    return {"msg": f"欢迎, {payload['sub']}"}

# 获取用户信息接口
@app.get("/auth/getUserInfo")
async def get_user_info(payload: dict = Depends(verify_token)):
    try:
        user_id = int(payload.get("sub", payload.get("userId", 0)))
        if not user_id:
            raise HTTPException(401, "用户ID无效")
        
        async with SessionLocal() as session:
            result = await session.execute(select(User).where(User.id == user_id))
            user = result.scalar_one_or_none()
            if not user:
                print(f"[getUserInfo] 用户不存在: user_id={user_id}")
                raise HTTPException(404, "用户不存在")
            
            # 安全地获取用户信息，处理可能为 None 的字段
            return {
                "code": 0,
                "data": {
                    "userId": user.id,
                    "userName": user.username or "",
                    "email": user.contact or "",
                    "realName": user.real_name or "",
                    "nickName": user.nick_name or "",
                    "theme": user.theme or "auto",
                    "language": user.language or "zh",
                    "avatar": user.avatar or "",
                    "department": user.department or "",
                    "roles": ["admin"],
                    "menus": [
                        {
                            "path": "/workbench",
                            "name": "工作台",
                            "icon": "icon-workbench"
                        },
                        {
                            "path": "/schedule",
                            "name": "日程管理",
                            "icon": "icon-calendar"
                        },
                        {
                            "path": "/meetings",
                            "name": "会议管理",
                            "icon": "icon-meeting"
                        },
                        {
                            "path": "/knowledge",
                            "name": "知识库",
                            "icon": "icon-knowledge"
                        },
                        {
                            "path": "/assistant",
                            "name": "智能助手",
                            "icon": "icon-ai"
                        },
                        {
                            "path": "/system/user-center",
                            "name": "个人中心",
                            "icon": "icon-user"
                        }
                    ],
                    "permissions": ["view_dashboard", "edit_profile", "manage_schedule", "manage_meetings", "use_ai_assistant"]
                }
            }
    except HTTPException:
        raise
    except ValueError as e:
        print(f"[getUserInfo] 用户ID转换错误: {e}, payload: {payload}")
        raise HTTPException(400, "用户ID格式错误")
    except Exception as e:
        print(f"[getUserInfo] 获取用户信息失败: {e}")
        import traceback
        traceback.print_exc()
        raise HTTPException(500, f"获取用户信息失败: {str(e)}")

# 初始化数据库
async def init_db():
    async with engine.begin() as conn:
        await conn.run_sync(Base.metadata.create_all)
        
    # 检查并添加新字段到today_tasks表
    try:
        async with SessionLocal() as session:
            # 检查表是否存在
            try:
                await session.execute(text("SELECT name FROM sqlite_master WHERE type='table' AND name='today_tasks'"))
                table_exists = True
            except Exception:
                table_exists = False
            
            if table_exists:
                # 尝试查询user_id字段，如果出错则添加字段
                try:
                    await session.execute(text("SELECT user_id FROM today_tasks LIMIT 1"))
                    print("today_tasks 表已包含 user_id 字段")
                except Exception:
                    # 添加user_id字段
                    try:
                        await session.execute(text("ALTER TABLE today_tasks ADD COLUMN user_id INTEGER"))
                        await session.commit()
                        print("已添加 user_id 字段到 today_tasks 表")
                        
                        # 为现有任务分配默认用户ID（使用第一个用户，如果没有用户则设为0）
                        try:
                            result = await session.execute(select(User).limit(1))
                            first_user = result.scalar_one_or_none()
                            default_user_id = first_user.id if first_user else 0
                            
                            if default_user_id > 0:
                                await session.execute(
                                    text("UPDATE today_tasks SET user_id = :user_id WHERE user_id IS NULL"),
                                    {"user_id": default_user_id}
                                )
                                await session.commit()
                                print(f"已为现有任务分配默认用户ID: {default_user_id}")
                            else:
                                # 如果没有用户，删除所有现有任务
                                await session.execute(text("DELETE FROM today_tasks WHERE user_id IS NULL"))
                                await session.commit()
                                print("已删除无用户关联的任务")
                        except Exception as e:
                            print(f"迁移现有任务数据失败: {e}")
                    except Exception as alter_error:
                        print(f"添加 user_id 字段失败: {alter_error}")
            else:
                print("today_tasks 表不存在，将在创建表时自动包含 user_id 字段")
    except Exception as e:
        print(f"today_tasks 表字段检查/添加失败: {e}")
        import traceback
        traceback.print_exc()
        
    # 检查并添加新字段到meetings表
    try:
        async with MeetingSessionLocal() as session:
            # 尝试查询新字段，如果出错则添加字段
            try:
                await session.execute(text("SELECT participants FROM meetings LIMIT 1"))
            except Exception:
                # 添加participants字段
                await session.execute(text("ALTER TABLE meetings ADD COLUMN participants TEXT DEFAULT ''"))
                await session.commit()
                print("已添加 participants 字段到 meetings 表")
            
            try:
                await session.execute(text("SELECT host_user_id FROM meetings LIMIT 1"))
            except Exception:
                # 添加host_user_id字段
                await session.execute(text("ALTER TABLE meetings ADD COLUMN host_user_id INTEGER DEFAULT 0"))
                await session.commit()
                print("已添加 host_user_id 字段到 meetings 表")
    except Exception as e:
        print(f"数据库字段检查/添加失败: {e}")
    
    async with meeting_engine.begin() as conn:
        await conn.run_sync(Base.metadata.create_all)

@app.on_event("startup")
async def on_startup():
    await init_db()

# 测试路由
@app.get("/")
async def root():
    return {"message": "API服务运行正常"}

# 健康检查接口
@app.get("/health")
async def health_check():
    return {"status": "ok"}

app.include_router(search.router)
app.include_router(contact.router)
app.include_router(group.router)
app.include_router(message.router)

# ========== 智能文档相关接口，操作 doc.json ==========
@app.post("/api/doc/upload")
async def upload_doc(file: UploadFile = File(...), payload: dict = Depends(verify_token)):
    user_id = payload["userId"]
    uploads_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'uploads')
    doc_path = os.path.join(uploads_dir, 'doc.json')
    os.makedirs(uploads_dir, exist_ok=True)
    save_path = os.path.join(uploads_dir, file.filename)
    with open(save_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
    # 生成自增数字id
    if os.path.exists(doc_path):
        with open(doc_path, 'r', encoding='utf-8') as f:
            try:
                doc_list = json.load(f)
            except Exception:
                doc_list = []
    else:
        doc_list = []
    # 检查是否已存在同名文件（仅当前用户）
    exist_id = None
    for item in doc_list:
        if item.get("filename") == file.filename and item.get("userId") == user_id:
            exist_id = item.get("id")
            break
    if exist_id is not None:
        doc_id = exist_id
    else:
        max_id = 0
        for item in doc_list:
            if item.get("userId") == user_id:
                try:
                    iid = int(item.get("id", 0))
                    if iid > max_id:
                        max_id = iid
                except Exception:
                    pass
        doc_id = max_id + 1
        # 自动读取内容
        ext = os.path.splitext(file.filename)[-1].lower()
        content = ""
        try:
            if ext == ".txt":
                with open(save_path, "r", encoding="utf-8", errors="ignore") as f:
                    content = f.read()
            elif ext == ".docx":
                from docx import Document
                doc = Document(save_path)
                content = "\n".join([p.text for p in doc.paragraphs])
            elif ext == ".doc":
                # 处理旧版Word格式(.doc)
                try:
                    # 尝试使用textract库（如果已安装）
                    import textract
                    content = textract.process(save_path).decode('utf-8')
                except ImportError:
                    # 如果textract未安装，尝试使用pypandoc
                    try:
                        import pypandoc
                        content = pypandoc.convert_file(save_path, 'plain', format='doc')
                    except ImportError:
                        # 如果都未安装，提示用户安装
                        content = "[错误: 需要安装textract或pypandoc库来解析.doc文件。请运行: pip install textract 或 pip install pypandoc]"
                    except Exception as e:
                        content = f"[解析DOC文件时出错: {str(e)}]"
                except Exception as e:
                    content = f"[解析DOC文件时出错: {str(e)}]"
            elif ext == ".pptx":
                from pptx import Presentation
                prs = Presentation(save_path)
                slides = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            slides.append(shape.text)
                content = "\n".join(slides)
            elif ext in [".xls", ".xlsx"]:
                try:
                    import openpyxl
                    wb = openpyxl.load_workbook(save_path)
                    sheets_content = []
                    for sheet_name in wb.sheetnames:
                        sheet = wb[sheet_name]
                        rows = []
                        for row in sheet.iter_rows(values_only=True):
                            row_data = [str(cell) if cell is not None else "" for cell in row]
                            rows.append("\t".join(row_data))
                        sheets_content.append(f"工作表: {sheet_name}\n" + "\n".join(rows))
                    content = "\n\n---工作表分隔---\n\n".join(sheets_content)
                except ImportError:
                    content = "[错误: 需要安装openpyxl库]"
                except Exception as e:
                    content = f"[解析Excel文件时出错: {str(e)}]"
            elif ext == ".csv":
                try:
                    import csv
                    with open(save_path, 'r', encoding='utf-8', errors='ignore') as f:
                        reader = csv.reader(f)
                        rows = []
                        for row in reader:
                            rows.append("\t".join(row))
                        content = "\n".join(rows)
                except Exception as e:
                    content = f"[解析CSV文件时出错: {str(e)}]"
            elif ext == ".pdf":
                try:
                    import PyPDF2
                    with open(save_path, 'rb') as f:
                        pdf_reader = PyPDF2.PdfReader(f)
                        pages_text = []
                        for page_num in range(len(pdf_reader.pages)):
                            page = pdf_reader.pages[page_num]
                            pages_text.append(page.extract_text())
                        content = "\n\n".join(pages_text)
                except ImportError:
                    content = "[错误: 需要安装PyPDF2库]"
                except Exception as e:
                    content = f"[解析PDF文件时出错: {str(e)}]"
        except Exception:
            content = ""
        import datetime
        upload_time = datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')
        doc_list.append({"id": doc_id, "filename": file.filename, "content": content, "upload_time": upload_time, "userId": user_id})
        with open(doc_path, 'w', encoding='utf-8') as f:
            json.dump(doc_list, f, ensure_ascii=False, indent=2)
    return {"id": doc_id, "filename": file.filename, "msg": "上传成功"}

@app.get("/api/doc/search")
async def search_docs(q: str = Query("", alias="q"), page: int = 1, size: int = 20, payload: dict = Depends(verify_token)):
    user_id = payload["userId"]
    uploads_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'uploads')
    doc_path = os.path.join(uploads_dir, 'doc.json')
    if not os.path.exists(doc_path):
        return {"total": 0, "data": []}
    with open(doc_path, 'r', encoding='utf-8') as f:
        doc_list = json.load(f)
    doc_list = [item for item in doc_list if item.get("userId") == user_id]
    if q:
        doc_list = [item for item in doc_list if q.lower() in item.get("filename", "").lower() or q.lower() in item.get("content", "").lower()]
    total = len(doc_list)
    start = (page - 1) * size
    end = start + size
    return {"total": total, "data": doc_list[start:end]}

@app.delete("/api/doc/{doc_id}")
async def delete_doc(doc_id: int = Path(...), payload: dict = Depends(verify_token)):
    user_id = payload["userId"]
    uploads_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'uploads')
    doc_path = os.path.join(uploads_dir, 'doc.json')
    if not os.path.exists(doc_path):
        return {"code": 1, "msg": "文档不存在"}
    with open(doc_path, 'r', encoding='utf-8') as f:
        doc_list = json.load(f)
    
    # 检查权限：只能删除自己的文档
    doc_to_delete = None
    for item in doc_list:
        if str(item.get('id')) == str(doc_id):
            doc_to_delete = item
            break
    
    if not doc_to_delete:
        return {"code": 1, "msg": "文档不存在"}
    
    if doc_to_delete.get("userId") != user_id:
        return {"code": 1, "msg": "无权限删除此文档"}
    
    new_list = [item for item in doc_list if str(item.get('id')) != str(doc_id)]
    with open(doc_path, 'w', encoding='utf-8') as f:
        json.dump(new_list, f, ensure_ascii=False, indent=2)
    return {"code": 0, "msg": "删除成功"}

@app.put("/api/doc/{doc_id}")
async def edit_doc(doc_id: int = Path(...), data: dict = Body(...), payload: dict = Depends(verify_token)):
    user_id = payload["userId"]
    uploads_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'uploads')
    doc_path = os.path.join(uploads_dir, 'doc.json')
    if not os.path.exists(doc_path):
        return {"code": 1, "msg": "文档不存在"}
    with open(doc_path, 'r', encoding='utf-8') as f:
        doc_list = json.load(f)
    updated = False
    for item in doc_list:
        if str(item.get('id')) == str(doc_id):
            # 检查权限：只能编辑自己的文档
            if item.get("userId") != user_id:
                return {"code": 1, "msg": "无权限编辑此文档"}
            item['content'] = data.get('html_content', item.get('content', ''))
            updated = True
            break
    if not updated:
        return {"code": 1, "msg": "文档不存在"}
    with open(doc_path, 'w', encoding='utf-8') as f:
        json.dump(doc_list, f, ensure_ascii=False, indent=2)
    return {"code": 0, "msg": "保存成功"}

@app.get("/api/doc/{doc_id}/download")
async def download_doc(doc_id: int = Path(...), payload: dict = Depends(verify_token)):
    user_id = payload["userId"]
    uploads_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'uploads')
    doc_path = os.path.join(uploads_dir, 'doc.json')
    if not os.path.exists(doc_path):
        return {"code": 1, "msg": "文档不存在"}
    with open(doc_path, 'r', encoding='utf-8') as f:
        doc_list = json.load(f)
    filename = None
    doc_user_id = None
    for item in doc_list:
        if str(item.get('id')) == str(doc_id):
            filename = item.get('filename')
            doc_user_id = item.get('userId')
            break
    if not filename:
        return {"code": 1, "msg": "文档不存在"}
    
    # 检查权限：只能下载自己的文档
    if doc_user_id != user_id:
        return {"code": 1, "msg": "无权限下载此文档"}
    
    file_path = os.path.join(uploads_dir, filename)
    
    # 如果原始文件存在，直接下载
    if os.path.exists(file_path):
        # 检测文件类型
        media_type, _ = mimetypes.guess_type(file_path)
        if not media_type:
            media_type = 'application/octet-stream'
        
        # 特殊处理常见文件类型
        file_ext = os.path.splitext(filename)[1].lower()
        if file_ext == '.pdf':
            media_type = 'application/pdf'
        elif file_ext in ['.doc', '.docx']:
            media_type = 'application/vnd.openxmlformats-officedocument.wordprocessingml.document' if file_ext == '.docx' else 'application/msword'
        elif file_ext in ['.xls', '.xlsx']:
            media_type = 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet' if file_ext == '.xlsx' else 'application/vnd.ms-excel'
        elif file_ext in ['.ppt', '.pptx']:
            media_type = 'application/vnd.openxmlformats-officedocument.presentationml.presentation' if file_ext == '.pptx' else 'application/vnd.ms-powerpoint'
        elif file_ext == '.txt':
            media_type = 'text/plain; charset=utf-8'
        
        # 处理中文文件名编码
        from urllib.parse import quote
        encoded_filename = quote(filename.encode('utf-8'))
        
        response = FileResponse(
            file_path,
            media_type=media_type
        )
        # 设置Content-Disposition头，使用RFC 5987格式支持中文文件名
        response.headers['Content-Disposition'] = f"attachment; filename*=UTF-8''{encoded_filename}"
        return response
    
    # 如果原始文件不存在，导出文档内容为文本文件
    doc_content = None
    for item in doc_list:
        if str(item.get('id')) == str(doc_id):
            doc_content = item.get('content', '')
            break
    
    if not doc_content:
        return {"code": 1, "msg": "文档内容不存在"}
    
    # 移除HTML标签（如果存在）
    if doc_content and '<' in doc_content and '>' in doc_content:
        import re
        import html
        doc_content = re.sub(r'<[^>]+>', '', doc_content)
        # 解码HTML实体
        doc_content = html.unescape(doc_content)
    
    # 创建临时文件
    temp_dir = os.path.join(uploads_dir, 'temp')
    os.makedirs(temp_dir, exist_ok=True)
    temp_file_path = os.path.join(temp_dir, f'doc_{doc_id}_{int(time.time())}.txt')
    
    with open(temp_file_path, 'w', encoding='utf-8') as f:
        f.write(doc_content)
    
    # 处理文件名（确保有.txt扩展名）
    download_filename = filename
    if '.' not in download_filename:
        download_filename += '.txt'
    elif not download_filename.lower().endswith('.txt'):
        # 如果原文件名有其他扩展名，改为.txt
        download_filename = os.path.splitext(download_filename)[0] + '.txt'
    
    # 处理中文文件名编码
    from urllib.parse import quote
    encoded_filename = quote(download_filename.encode('utf-8'))
    
    response = FileResponse(
        temp_file_path,
        media_type='text/plain; charset=utf-8'
    )
    # 设置Content-Disposition头，使用RFC 5987格式支持中文文件名
    response.headers['Content-Disposition'] = f"attachment; filename*=UTF-8''{encoded_filename}"
    
    return response

@app.get('/api/doc/{doc_id}')
async def get_doc_detail(doc_id: int, payload: dict = Depends(verify_token)):
    user_id = payload["userId"]
    uploads_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'uploads')
    doc_path = os.path.join(uploads_dir, 'doc.json')
    if not os.path.exists(doc_path):
        return {"code": 1, "msg": "文档不存在"}
    with open(doc_path, 'r', encoding='utf-8') as f:
        doc_list = json.load(f)
    for item in doc_list:
        if str(item.get('id')) == str(doc_id):
            # 检查权限：只能查看自己的文档
            if item.get("userId") != user_id:
                return {"code": 1, "msg": "无权限访问此文档"}
            return {"code": 0, "content": item.get('content', '')}
    return {"code": 1, "msg": "未找到文档"}

# ========== 知识库相关接口，操作 knowledge.json ==========
@app.get("/api/knowledge")
async def get_knowledge(q: str = Query("", alias="q"), page: int = 1, size: int = 20, payload: dict = Depends(verify_token)):
    # 移除用户过滤，所有用户都可以查看所有知识库数据
    uploads_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'uploads')
    kb_path = os.path.join(uploads_dir, 'knowledge.json')
    if not os.path.exists(kb_path):
        return {"total": 0, "data": []}
    try:
        with open(kb_path, 'r', encoding='utf-8') as f:
            kb_list = json.load(f)
    except Exception as e:
        print(f'知识库读取失败: {e}')
        kb_list = []
    # 移除用户过滤，显示所有文档
    if q:
        kb_list = [item for item in kb_list if q.lower() in item.get("filename", "").lower() or q.lower() in item.get("content", "").lower()]
    
    # 按ID排序（升序）
    kb_list.sort(key=lambda x: int(x.get("id", 0)) if x.get("id") is not None else 0)
    
    total = len(kb_list)
    start = (page - 1) * size
    end = start + size
    return {"total": total, "data": kb_list[start:end]}

@app.get("/api/knowledge/{doc_id}/download")
async def download_knowledge_doc(doc_id: int = Path(...), payload: dict = Depends(verify_token)):
    """下载知识库文档"""
    user_id = payload["userId"]
    uploads_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'uploads')
    kb_path = os.path.join(uploads_dir, 'knowledge.json')
    doc_path = os.path.join(uploads_dir, 'doc.json')
    
    if not os.path.exists(kb_path):
        return {"code": 1, "msg": "知识库不存在"}
    
    try:
        with open(kb_path, 'r', encoding='utf-8') as f:
            kb_list = json.load(f)
        
        # 查找文档（同时检查docId和id）
        doc_item = None
        for item in kb_list:
            item_doc_id = item.get('docId')
            item_id = item.get('id')
            if (item_doc_id is not None and str(item_doc_id) == str(doc_id)) or \
               (item_id is not None and str(item_id) == str(doc_id)):
                doc_item = item
                break
        
        if not doc_item:
            return JSONResponse({"code": 1, "msg": "文档不存在"}, status_code=404)
        
        # 优先尝试从智能文档中下载原始文件
        doc_id_for_file = doc_item.get('docId')
        if doc_id_for_file and os.path.exists(doc_path):
            with open(doc_path, 'r', encoding='utf-8') as f:
                doc_list = json.load(f)
            for doc in doc_list:
                if str(doc.get('id')) == str(doc_id_for_file):
                    filename = doc.get('filename')
                    if filename:
                        file_path = os.path.join(uploads_dir, filename)
                        if os.path.exists(file_path):
                            # 正确处理中文文件名编码
                            from urllib.parse import quote
                            # 对文件名进行URL编码，支持中文（RFC 5987格式）
                            encoded_filename = quote(filename.encode('utf-8'))
                            
                            # 根据文件扩展名检测MIME类型
                            file_ext = os.path.splitext(filename)[1].lower()
                            media_type = mimetypes.guess_type(filename)[0] or 'application/octet-stream'
                            
                            # 特殊处理常见文件类型
                            if file_ext == '.pdf':
                                media_type = 'application/pdf'
                            elif file_ext in ['.doc', '.docx']:
                                media_type = 'application/vnd.openxmlformats-officedocument.wordprocessingml.document' if file_ext == '.docx' else 'application/msword'
                            elif file_ext in ['.xls', '.xlsx']:
                                media_type = 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet' if file_ext == '.xlsx' else 'application/vnd.ms-excel'
                            elif file_ext in ['.ppt', '.pptx']:
                                media_type = 'application/vnd.openxmlformats-officedocument.presentationml.presentation' if file_ext == '.pptx' else 'application/vnd.ms-powerpoint'
                            elif file_ext == '.txt':
                                media_type = 'text/plain; charset=utf-8'
                            
                            response = FileResponse(
                                file_path,
                                media_type=media_type
                            )
                            # 设置Content-Disposition头，使用RFC 5987格式支持中文文件名
                            response.headers['Content-Disposition'] = f"attachment; filename*=UTF-8''{encoded_filename}"
                            return response
        
        # 如果没有原始文件，将内容导出为文本文件
        filename = doc_item.get('filename') or doc_item.get('title') or f'文档_{doc_id}'
        # 确保文件名有扩展名
        if '.' not in filename:
            filename += '.txt'
        
        # 获取文档内容
        content = doc_item.get('content', '')
        # 移除HTML标签（如果存在）
        if content and '<' in content and '>' in content:
            import re
            content = re.sub(r'<[^>]+>', '', content)
            # 解码HTML实体
            content = html.unescape(content)
        
        # 创建临时文件
        temp_dir = os.path.join(uploads_dir, 'temp')
        os.makedirs(temp_dir, exist_ok=True)
        temp_file_path = os.path.join(temp_dir, f'kb_doc_{doc_id}_{int(time.time())}.txt')
        
        with open(temp_file_path, 'w', encoding='utf-8') as f:
            f.write(content)
        
        # 返回文件，正确处理中文文件名编码
        from urllib.parse import quote
        # 对文件名进行URL编码，支持中文（RFC 5987格式）
        encoded_filename = quote(filename.encode('utf-8'))
        
        response = FileResponse(
            temp_file_path, 
            media_type='text/plain; charset=utf-8'
        )
        # 设置Content-Disposition头，使用RFC 5987格式支持中文文件名
        # 使用filename*参数，这是RFC 5987标准，支持UTF-8编码
        response.headers['Content-Disposition'] = f"attachment; filename*=UTF-8''{encoded_filename}"
        return response
    
    except Exception as e:
        import traceback
        traceback.print_exc()
        return JSONResponse({"code": 1, "msg": f"下载失败: {str(e)}"}, status_code=500)

@app.get("/api/knowledge/detail")
async def knowledge_detail(id: int = Query(...), payload: dict = Depends(verify_token)):
    # 移除权限检查，所有用户都可以查看文档详情
    uploads_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'uploads')
    kb_path = os.path.join(uploads_dir, 'knowledge.json')
    if not os.path.exists(kb_path):
        return {"code": 1, "msg": "知识库不存在"}
    with open(kb_path, 'r', encoding='utf-8') as f:
        kb_list = json.load(f)
    for item in kb_list:
        if str(item.get('docId', item.get('id'))) == str(id):
            # 移除权限检查，所有用户都可以访问
            return {"code": 0, "data": item}
    return {"code": 1, "msg": "未找到文档"}

@app.post("/api/knowledge/edit")
async def knowledge_edit(data: dict = Body(...), payload: dict = Depends(verify_token)):
    user_id = payload["userId"]
    doc_id = data.get("id")
    content = data.get("content")
    knowledge_base = data.get("knowledge_base")  # 支持更新知识库
    uploads_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'uploads')
    kb_path = os.path.join(uploads_dir, 'knowledge.json')
    if not os.path.exists(kb_path):
        return {"code": 1, "msg": "知识库不存在"}
    with open(kb_path, 'r', encoding='utf-8') as f:
        kb_list = json.load(f)
    updated = False
    for item in kb_list:
        if str(item.get('docId', item.get('id'))) == str(doc_id):
            # 检查文档所有权
            if item.get("userId") != user_id:
                return {"code": 1, "msg": "无权限编辑此文档"}
            if content is not None:
                item['content'] = content
            if knowledge_base is not None:
                item['knowledge_base'] = knowledge_base
            updated = True
            break
    if not updated:
        return {"code": 1, "msg": "文档不存在"}
    with open(kb_path, 'w', encoding='utf-8') as f:
        json.dump(kb_list, f, ensure_ascii=False, indent=2)
    return {"code": 0, "msg": "保存成功"}

@app.post("/api/knowledge/upload")
async def knowledge_upload(
    file: UploadFile = File(...), 
    knowledge_base: str = Form(None),
    payload: dict = Depends(verify_token)
):
    user_id = payload["userId"]
    uploads_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'uploads')
    kb_path = os.path.join(uploads_dir, 'knowledge.json')
    os.makedirs(uploads_dir, exist_ok=True)
    # 检查是否已存在同名文档（仅当前用户）
    if os.path.exists(kb_path):
        with open(kb_path, 'r', encoding='utf-8') as f:
            try:
                kb_list = json.load(f)
            except Exception:
                kb_list = []
        for item in kb_list:
            if item.get("filename") == file.filename and item.get("userId") == user_id:
                return {"code": 1, "msg": "该文档已存在，禁止重复上传"}
    else:
        kb_list = []
    save_path = os.path.join(uploads_dir, file.filename)
    with open(save_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
    # 自动读取内容
    ext = os.path.splitext(file.filename)[-1].lower()
    content = ""
    try:
        if ext == ".txt":
            with open(save_path, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read()
        elif ext == ".docx":
            from docx import Document
            doc = Document(save_path)
            content = "\n".join([p.text for p in doc.paragraphs])
        elif ext == ".doc":
            # 处理旧版Word格式(.doc)
            try:
                # 尝试使用textract库（如果已安装）
                import textract
                content = textract.process(save_path).decode('utf-8')
            except ImportError:
                # 如果textract未安装，尝试使用pypandoc
                try:
                    import pypandoc
                    content = pypandoc.convert_file(save_path, 'plain', format='doc')
                except ImportError:
                    # 如果都未安装，提示用户安装
                    content = "[错误: 需要安装textract或pypandoc库来解析.doc文件。请运行: pip install textract 或 pip install pypandoc]"
                except Exception as e:
                    content = f"[解析DOC文件时出错: {str(e)}]"
            except Exception as e:
                content = f"[解析DOC文件时出错: {str(e)}]"
        elif ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(save_path)
            slides = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slides.append(shape.text)
            content = "\n".join(slides)
        elif ext == ".pdf":
            try:
                import PyPDF2
                with open(save_path, 'rb') as f:
                    pdf_reader = PyPDF2.PdfReader(f)
                    pages_text = []
                    for page_num in range(len(pdf_reader.pages)):
                        page = pdf_reader.pages[page_num]
                        pages_text.append(page.extract_text())
                    content = "\n\n".join(pages_text)
            except ImportError:
                content = "[错误: 需要安装PyPDF2库]"
            except Exception as e:
                content = f"[解析PDF文件时出错: {str(e)}]"
    except Exception:
        content = ""
    import datetime
    upload_time = datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')
    # 生成自增id - 检查所有文档的ID，确保全局唯一
    max_id = 0
    for item in kb_list:
        # 移除用户过滤，检查所有文档的ID以确保全局唯一
            try:
                iid = int(item.get("id", 0))
                if iid > max_id:
                    max_id = iid
            except Exception:
                pass
    doc_id = max_id + 1
    kb_list.append({
        "id": doc_id,
        "filename": file.filename,
        "content": content,
        "upload_time": upload_time,
        "userId": user_id,
        "knowledge_base": knowledge_base or ""  # 保存知识库分类
    })
    with open(kb_path, 'w', encoding='utf-8') as f:
        json.dump(kb_list, f, ensure_ascii=False, indent=2)
    return {"code": 0, "msg": "上传成功", "filename": file.filename, "id": doc_id}

@app.post("/api/knowledge/delete")
async def knowledge_delete(data: dict = Body(...), payload: dict = Depends(verify_token)):
    user_id = payload["userId"]
    ids = data.get("ids")
    if not ids:
        return {"code": 1, "msg": "未指定要删除的文档id"}
    uploads_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'uploads')
    kb_path = os.path.join(uploads_dir, 'knowledge.json')
    if not os.path.exists(kb_path):
        return {"code": 1, "msg": "知识库不存在"}
    with open(kb_path, 'r', encoding='utf-8') as f:
        kb_list = json.load(f)
    
    # 检查权限：只能删除属于自己的文档
    ids_to_delete = [str(i) for i in ids]
    new_list = []
    deleted_count = 0
    
    for item in kb_list:
        item_doc_id = item.get('docId')
        item_id = item.get('id')
        # 检查是否在删除列表中（同时检查docId和id）
        should_delete = False
        for delete_id in ids_to_delete:
            if (item_doc_id is not None and str(item_doc_id) == str(delete_id)) or \
               (item_id is not None and str(item_id) == str(delete_id)):
                should_delete = True
                break
        
        if should_delete:
            # 获取当前用户信息，检查部门和角色
            current_user_department = ""
            async with SessionLocal() as session:
                result = await session.execute(select(User).where(User.id == user_id))
                current_user = result.scalar_one_or_none()
                if current_user:
                    current_user_department = current_user.department or ""
            
            # 定义有删除所有文档权限的部门
            admin_departments = ["人事部", "技术部", "管理部", "行政部"]
            is_admin_user = current_user_department in admin_departments
            
            # 权限检查：1. 文档所有者可以删除 2. 管理员部门可以删除所有文档
            if item.get("userId") == user_id or is_admin_user:
                deleted_count += 1
                continue  # 跳过此项（即删除）
            else:
                # 无权限删除，保留此项
                new_list.append(item)
        else:
            # 不在删除列表中，保留此项
            new_list.append(item)
    
    if deleted_count == 0:
        return {"code": 1, "msg": "没有可删除的文档或无权限删除", "deleted_count": 0}
    
    with open(kb_path, 'w', encoding='utf-8') as f:
        json.dump(new_list, f, ensure_ascii=False, indent=2)
    return {"code": 0, "msg": f"成功删除 {deleted_count} 个文档", "deleted_count": deleted_count}

@app.get("/api/knowledge/search")
async def knowledge_search(kw: str = Query(""), payload: dict = Depends(verify_token)):
    # 移除用户过滤，所有用户都可以搜索所有文档
    uploads_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'uploads')
    kb_path = os.path.join(uploads_dir, 'knowledge.json')
    if not os.path.exists(kb_path):
        return {"data": []}
    with open(kb_path, 'r', encoding='utf-8') as f:
        kb_list = json.load(f)
    
    # 移除用户过滤，搜索所有文档
    result = [
        item for item in kb_list
        if kw.lower() in (item.get("title", "") + item.get("filename", "")).lower()
        or kw.lower() in item.get("content", "").lower()
    ]
    # 按ID排序（升序）
    result.sort(key=lambda x: int(x.get("id", 0)) if x.get("id") is not None else 0)
    return {"data": result}

@app.post("/api/knowledge/restore")
async def restore_knowledge_docs(payload: dict = Depends(verify_token)):
    """从备份文件或智能文档中恢复被覆盖删除的文档"""
    import shutil
    from datetime import datetime
    import glob
    
    user_id = payload["userId"]
    uploads_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'uploads')
    kb_path = os.path.join(uploads_dir, 'knowledge.json')
    doc_path = os.path.join(uploads_dir, 'doc.json')
    
    try:
        # 读取当前知识库
        current_kb_list = []
        if os.path.exists(kb_path):
            with open(kb_path, 'r', encoding='utf-8') as f:
                current_kb_list = json.load(f)
        
        # 读取智能文档列表
        doc_list = []
        if os.path.exists(doc_path):
            with open(doc_path, 'r', encoding='utf-8') as f:
                doc_list = json.load(f)
        
        # 获取所有备份文件，按时间排序（最新的在前）
        backup_files = glob.glob(os.path.join(uploads_dir, 'knowledge.json.backup.*'))
        backup_files.sort(reverse=True)
        
        restored_count = 0
        restored_docs = []
        
        # 方法1：从最新的备份文件中恢复丢失的文档
        if backup_files:
            latest_backup = backup_files[0]
            with open(latest_backup, 'r', encoding='utf-8') as f:
                backup_kb_list = json.load(f)
            
            # 获取当前知识库中所有docId的集合
            current_doc_ids = {str(item.get('docId', item.get('id'))) for item in current_kb_list}
            
            # 从备份中找到丢失的文档
            for backup_item in backup_kb_list:
                backup_doc_id = str(backup_item.get('docId', backup_item.get('id')))
                # 如果备份中的文档不在当前知识库中，则恢复
                if backup_doc_id not in current_doc_ids:
                    # 检查是否已存在相同ID的文档（避免重复）
                    exists = False
                    for current_item in current_kb_list:
                        if str(current_item.get('id')) == str(backup_item.get('id')):
                            exists = True
                            break
                    
                    if not exists:
                        current_kb_list.append(backup_item)
                        restored_count += 1
                        restored_docs.append({
                            'id': backup_item.get('id'),
                            'title': backup_item.get('title', backup_item.get('filename', ''))
                        })
        
        # 方法2：从智能文档中恢复应该存在但丢失的文档
        # 获取知识库中所有docId
        kb_doc_ids = {str(item.get('docId')) for item in current_kb_list if item.get('docId')}
        
        # 检查智能文档中哪些应该被推送到知识库但不存在
        for doc_item in doc_list:
            doc_id = str(doc_item.get('id'))
            # 如果智能文档不在知识库中，且是当前用户的文档，则恢复
            if doc_id not in kb_doc_ids and str(doc_item.get('userId')) == str(user_id):
                # 检查是否已存在相同ID的知识库文档
                exists = False
                for kb_item in current_kb_list:
                    if str(kb_item.get('id')) == doc_id or str(kb_item.get('docId')) == doc_id:
                        exists = True
                        break
                
                if not exists:
                    # 重新生成知识库ID，确保不冲突
                    max_id = 0
                    used_ids = set()
                    for item in current_kb_list:
                        try:
                            kb_id = int(item.get("id", 0))
                            if kb_id > max_id:
                                max_id = kb_id
                            if kb_id > 0:
                                used_ids.add(kb_id)
                        except Exception:
                            pass
                    
                    new_kb_id = max_id + 1
                    while new_kb_id == int(doc_id) or new_kb_id in used_ids:
                        new_kb_id += 1
                    
                    # 添加恢复的文档
                    restored_item = {
                        "id": new_kb_id,
                        "docId": int(doc_id),
                        "title": doc_item.get('filename', ''),
                        "filename": doc_item.get('filename', ''),
                        "content": doc_item.get('content', ''),
                        "type": doc_item.get('type', '文档'),
                        "userId": user_id,
                        "upload_time": doc_item.get('upload_time', datetime.now().strftime('%Y-%m-%d %H:%M:%S')),
                        "knowledge_base": ""
                    }
                    current_kb_list.append(restored_item)
                    restored_count += 1
                    restored_docs.append({
                        'id': new_kb_id,
                        'title': restored_item.get('title', '')
                    })
        
        # 保存恢复后的知识库
        if restored_count > 0:
            # 创建备份
            backup_path = kb_path + f".backup.before_restore.{datetime.now().strftime('%Y%m%d_%H%M%S')}"
            if os.path.exists(kb_path):
                shutil.copy2(kb_path, backup_path)
            
            # 保存恢复后的数据
            with open(kb_path, 'w', encoding='utf-8') as f:
                json.dump(current_kb_list, f, ensure_ascii=False, indent=2)
            
            return {
                "code": 0,
                "msg": f"成功恢复 {restored_count} 个文档",
                "data": {
                    "restored_count": restored_count,
                    "restored_docs": restored_docs,
                    "backup_file": backup_path if os.path.exists(backup_path) else None
                }
            }
        else:
            return {
                "code": 0,
                "msg": "未发现需要恢复的文档",
                "data": {
                    "restored_count": 0,
                    "restored_docs": []
                }
            }
    
    except Exception as e:
        import traceback
        traceback.print_exc()
        return {"code": 1, "msg": f"恢复失败: {str(e)}"}

@app.post("/api/knowledge/fix-ids")
async def fix_knowledge_ids(payload: dict = Depends(verify_token)):
    """修复知识库中重复的ID，重新分配唯一ID"""
    import shutil
    from datetime import datetime
    
    uploads_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'uploads')
    kb_path = os.path.join(uploads_dir, 'knowledge.json')
    
    if not os.path.exists(kb_path):
        return {"code": 1, "msg": "知识库文件不存在"}
    
    try:
        # 读取原始数据
        with open(kb_path, 'r', encoding='utf-8') as f:
            kb_list = json.load(f)
        
        if not isinstance(kb_list, list):
            return {"code": 1, "msg": "数据格式不正确"}
        
        # 检测重复ID
        id_count = {}
        duplicates = []
        for idx, item in enumerate(kb_list):
            item_id = item.get('id')
            if item_id is None:
                continue
            if item_id in id_count:
                id_count[item_id].append(idx)
                duplicates.append(item_id)
            else:
                id_count[item_id] = [idx]
        
        # 创建备份
        backup_path = kb_path + f".backup.{datetime.now().strftime('%Y%m%d_%H%M%S')}"
        try:
            shutil.copy2(kb_path, backup_path)
        except Exception as e:
            return {"code": 1, "msg": f"创建备份失败: {e}"}
        
        # 重新分配ID（从1开始，每个文档都分配唯一ID）
        # 直接为每个文档分配新的唯一ID，确保每个文档都有唯一ID
        current_id = 1
        for item in kb_list:
            # 直接分配新ID，不管原来的ID是什么
            item['id'] = current_id
            # 同时更新docId字段，确保docId和id保持一致
            if 'docId' in item:
                item['docId'] = current_id
            current_id += 1
        
        # 验证新ID的唯一性
        new_ids = [item.get('id') for item in kb_list if item.get('id') is not None]
        if len(new_ids) != len(set(new_ids)):
            return {"code": 1, "msg": "重新分配后仍有重复ID"}
        
        # 保存更新后的数据
        with open(kb_path, 'w', encoding='utf-8') as f:
            json.dump(kb_list, f, ensure_ascii=False, indent=2)
        
        return {
            "code": 0,
            "msg": "修复成功",
            "data": {
                "total": len(kb_list),
                "unique_ids": len(set(new_ids)),
                "duplicates_fixed": len(set(duplicates)) if duplicates else 0,
                "backup_file": backup_path
            }
        }
    except Exception as e:
        import traceback
        traceback.print_exc()
        return {"code": 1, "msg": f"修复失败: {str(e)}"}

@app.post("/api/knowledge/qa")
async def knowledge_qa(data: dict = Body(...), payload: dict = Depends(verify_token)):
    """知识库智能问答接口 - 增强版"""
    user_id = payload["userId"]
    question = data.get("question", "")
    knowledge_base = data.get("knowledge_base", "")
    conversation_id = data.get("conversation_id", "")
    history = data.get("history", [])
    
    if not question.strip():
        return {"code": 1, "msg": "问题不能为空"}
    
    # 生成会话ID（如果没有）
    if not conversation_id:
        import uuid
        conversation_id = str(uuid.uuid4())
    
    # 读取知识库文档
    uploads_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'uploads')
    kb_path = os.path.join(uploads_dir, 'knowledge.json')
    
    if not os.path.exists(kb_path):
        return {"code": 0, "answer": "知识库为空，请先上传文档。"}
    
    try:
        with open(kb_path, 'r', encoding='utf-8') as f:
            kb_list = json.load(f)
    except Exception as e:
        return {"code": 1, "msg": f"读取知识库失败: {e}"}
    
    # 移除用户过滤，所有用户都可以使用所有文档进行问答
    # kb_list = [item for item in kb_list if item.get("userId") == user_id]  # 已移除用户过滤
    
    # 过滤知识库
    if knowledge_base:
        kb_list = [item for item in kb_list if item.get("knowledge_base") == knowledge_base]
    
    # 智能搜索相关文档
    question_lower = question.lower()
    relevant_docs = []
    
    # 意图识别
    intent_keywords = {
        "规定": ["规定", "制度", "政策", "规则", "条例"],
        "申请": ["申请", "如何", "流程", "步骤", "办理"],
        "预约": ["预约", "预定", "安排", "订"],
        "技术": ["开发", "技术", "代码", "API", "接口"],
        "培训": ["培训", "学习", "教程", "指南"],
        "联系": ["联系", "电话", "邮箱", "地址"]
    }
    
    # 根据历史对话优化搜索
    context_keywords = []
    if history:
        for item in history[-3:]:  # 最近3次对话
            context_keywords.extend(item.get("question", "").lower().split())
    
    for doc in kb_list:
        content = doc.get("content", "")
        title = doc.get("title", doc.get("filename", ""))
        
        # 计算匹配度
        score = 0
        
        # 1. 标题匹配（权重最高）
        title_lower = title.lower()
        for word in question_lower.split():
            if len(word) > 1 and word in title_lower:
                score += 15
        
        # 2. 意图匹配
        for intent, keywords in intent_keywords.items():
            if any(kw in question_lower for kw in keywords):
                for kw in keywords:
                    if kw in content.lower():
                        score += 8
        
        # 3. 内容匹配
        content_lower = content.lower()
        question_words = [w for w in question_lower.split() if len(w) > 1]
        for word in question_words:
            if word in content_lower:
                # 计算词频和位置权重
                word_count = content_lower.count(word)
                score += word_count * 2
                
                # 如果词汇出现在开头，加分
                if content_lower.find(word) < 100:
                    score += 3
        
        # 4. 上下文匹配
        for context_word in context_keywords:
            if len(context_word) > 1 and context_word in content_lower:
                score += 1
        
        # 5. 语义相关性（简单实现）
        semantic_score = calculate_semantic_similarity(question, content)
        score += semantic_score
        
        if score > 0:
            relevant_docs.append({
                "doc": doc,
                "score": score,
                "title": title,
                "content_preview": content[:200] + "..." if len(content) > 200 else content
            })
    
    # 按相关度排序
    relevant_docs.sort(key=lambda x: x["score"], reverse=True)
    
    if not relevant_docs:
        return {
            "code": 0, 
            "answer": f"抱歉，在知识库中没有找到与『{question}』相关的信息。建议您：<br/>1. 尝试使用更简单的关键词<br/>2. 检查是否有相关文档已上传<br/>3. 联系管理员获取更多帮助"
        }
    
    # 生成智能回答
    answer_parts = []
    answer_parts.append(f"基于知识库搜索，为您找到以下相关信息：")
    
    # 取前3个最相关的文档
    top_docs = relevant_docs[:3]
    
    for i, item in enumerate(top_docs, 1):
        doc = item["doc"]
        title = item["title"]
        content_preview = item["content_preview"]
        
        answer_parts.append(f"<br/><br/><strong>{i}. {title}</strong>")
        
        # 智能提取相关段落
        content = doc.get("content", "")
        relevant_sentences = []
        
        for word in question_lower.split():
            if len(word) > 1:
                sentences = content.split('。')
                for sentence in sentences:
                    if word in sentence.lower() and len(sentence.strip()) > 10:
                        relevant_sentences.append(sentence.strip() + '。')
                        if len(relevant_sentences) >= 2:
                            break
        
        if relevant_sentences:
            answer_parts.append("<br/>".join(relevant_sentences[:2]))
        else:
            answer_parts.append(content_preview)
    
    # 添加建议
    if len(relevant_docs) > 3:
        answer_parts.append(f"<br/><br/><small>💡 提示：还找到了 {len(relevant_docs) - 3} 个相关文档，您可以在文档列表中查看更多详细信息。</small>")
    
    answer_parts.append("<br/><br/><small>📚 如需查看完整文档内容，请在文档列表中点击预览按钮。</small>")
    
    return {
        "code": 0,
        "answer": "".join(answer_parts),
        "relevant_count": len(relevant_docs),
        "docs": [{"id": item["doc"].get("id"), "title": item["title"]} for item in top_docs],
        "conversation_id": conversation_id,
        "intent": detect_intent(question),
        "confidence": calculate_confidence(relevant_docs)
    }

def calculate_semantic_similarity(question: str, content: str) -> float:
    """计算语义相似性（简单实现）"""
    question_words = set(question.lower().split())
    content_words = set(content.lower().split())
    
    if not question_words or not content_words:
        return 0
    
    intersection = question_words.intersection(content_words)
    union = question_words.union(content_words)
    
    # Jaccard相似度
    similarity = len(intersection) / len(union) if union else 0
    return similarity * 10  # 放大到0-10分

def detect_intent(question: str) -> str:
    """检测问题意图"""
    question_lower = question.lower()
    
    intent_patterns = {
        "询问规定": ["规定", "制度", "政策", "规则", "条例"],
        "申请流程": ["申请", "如何", "流程", "步骤", "办理"],
        "预约服务": ["预约", "预定", "安排", "订"],
        "技术咨询": ["开发", "技术", "代码", "API", "接口"],
        "培训学习": ["培训", "学习", "教程", "指南"],
        "联系方式": ["联系", "电话", "邮箱", "地址"],
        "一般查询": []
    }
    
    for intent, keywords in intent_patterns.items():
        if any(kw in question_lower for kw in keywords):
            return intent
    
    return "一般查询"

def calculate_confidence(relevant_docs: list) -> float:
    """计算回答置信度"""
    if not relevant_docs:
        return 0.0
    
    max_score = relevant_docs[0]["score"]
    doc_count = len(relevant_docs)
    
    # 基于最高分数和文档数量计算置信度
    confidence = min(max_score / 50, 1.0)  # 最高分数归一化
    
    # 文档数量加分
    if doc_count >= 3:
        confidence += 0.1
    elif doc_count >= 2:
        confidence += 0.05
    
    return round(min(confidence, 1.0), 2)

@app.delete("/api/knowledge/{doc_id}")
async def delete_knowledge(doc_id: int = Path(...), payload: dict = Depends(verify_token)):
    user_id = payload["userId"]
    uploads_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'uploads')
    kb_path = os.path.join(uploads_dir, 'knowledge.json')
    if not os.path.exists(kb_path):
        return {"code": 1, "msg": "知识库不存在"}
    try:
        # 获取当前用户信息，检查部门和角色
        current_user_department = ""
        current_user_roles = []
        async with SessionLocal() as session:
            result = await session.execute(select(User).where(User.id == user_id))
            current_user = result.scalar_one_or_none()
            if current_user:
                current_user_department = current_user.department or ""
                # 可以从用户信息中获取角色，这里假设有roles字段或通过其他方式判断
                # 暂时通过部门判断权限
        
        # 定义有删除所有文档权限的部门
        admin_departments = ["人事部", "技术部", "管理部", "行政部"]
        is_admin_user = current_user_department in admin_departments
        
        with open(kb_path, 'r', encoding='utf-8') as f:
            kb_list = json.load(f)
        
        # 检查文档是否存在以及权限
        doc_found = False
        doc_owned = False
        doc_user_id = None
        for item in kb_list:
            if str(item.get('docId', item.get('id'))) == str(doc_id):
                doc_found = True
                doc_user_id = item.get("userId")
                if doc_user_id == user_id:
                    doc_owned = True
                break
        
        if not doc_found:
            return {"code": 1, "msg": "文档不存在"}
        
        # 权限检查：1. 文档所有者可以删除 2. 管理员部门可以删除所有文档
        if not doc_owned and not is_admin_user:
            return {"code": 1, "msg": f"无权限删除此文档。只有文档上传者或{', '.join(admin_departments)}的成员可以删除文档。"}
        
        # 删除文档（同时检查docId和id字段）
        new_list = []
        deleted = False
        for item in kb_list:
            item_doc_id = item.get('docId')
            item_id = item.get('id')
            # 匹配docId或id
            if str(item_doc_id) == str(doc_id) or str(item_id) == str(doc_id):
                deleted = True
                continue  # 跳过此项（即删除）
            new_list.append(item)
        
        if not deleted:
            return {"code": 1, "msg": "文档不存在或已被删除"}
        
        # 保存更新后的列表
        with open(kb_path, 'w', encoding='utf-8') as f:
            json.dump(new_list, f, ensure_ascii=False, indent=2)
        return {"code": 0, "msg": "删除成功"}
    except Exception as e:
        return {"code": 1, "msg": f"删除失败: {e}"}



# ========== 推送到知识库时，从 doc.json 读取内容 ==========
@app.post("/api/knowledge/import")
async def push_to_knowledge(data: dict = Body(...), payload: dict = Depends(verify_token)):
    user_id = payload["userId"]
    uploads_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'uploads')
    doc_path = os.path.join(uploads_dir, 'doc.json')
    kb_path = os.path.join(uploads_dir, 'knowledge.json')
    
    doc_id = data.get("docId")
    title = data.get("title")
    doc_type = data.get("type", "文档")
    content = data.get("content", "")
    
    # 如果没有传content，尝试从doc.json读取
    if not content and doc_id and os.path.exists(doc_path):
        with open(doc_path, 'r', encoding='utf-8') as f:
            doc_list = json.load(f)
        for item in doc_list:
            if str(item.get('id')) == str(doc_id):
                content = item.get('content', '')
                break
    
    # 如果内容仍然为空，创建基本信息
    if not content or content.strip() == '' or content.strip() == '<p><br></p>':
        content = f"<h1>{title}</h1><p>此文档已推送到知识库，但原始内容为空。</p><p>文档类型：{doc_type}</p><p>推送时间：{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}</p>"
    
    if not (doc_id and title):
        return {"code": 1, "msg": "缺少必要参数：文档ID或标题"}
    
    # 读取knowledge.json
    if os.path.exists(kb_path):
        with open(kb_path, 'r', encoding='utf-8') as f:
            try:
                kb_list = json.load(f)
            except Exception:
                kb_list = []
    else:
        kb_list = []
    
    
    # 查找是否已存在（同一用户的相同docId文档，避免重复推送）
    found = False
    for item in kb_list:
        if item.get("docId") == doc_id and item.get("userId") == user_id:
            # 如果已存在，更新现有记录（只更新内容，不改变ID）
            item["title"] = title
            item["filename"] = title
            item["content"] = content
            item["type"] = doc_type
            item["upload_time"] = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
            found = True
            break
    
    if not found:
        # 生成新的知识库ID - 检查所有文档的ID，确保全局唯一且不与智能文档ID冲突
        max_id = 0
        used_ids = set()
        for item in kb_list:
            try:
                kb_id = int(item.get("id", 0))
                if kb_id > max_id:
                    max_id = kb_id
                if kb_id > 0:
                    used_ids.add(kb_id)
            except Exception:
                pass
        
        # 生成新ID，确保不与现有ID重复，也不与智能文档ID相同
        new_id = max_id + 1
        # 如果新ID与doc_id相同，则继续递增直到找到不冲突的ID
        while new_id == int(doc_id) or new_id in used_ids:
            new_id += 1
        
        # 添加新记录，知识库ID与智能文档ID分开
        kb_list.append({
            "id": new_id,           # 知识库自增ID，与智能文档ID分开
            "docId": doc_id,       # 保存智能文档的ID作为关联字段
            "title": title,
            "filename": title,
            "content": content,
            "type": doc_type,
            "userId": user_id,
            "upload_time": datetime.now().strftime('%Y-%m-%d %H:%M:%S')
        })
    
    # 保存到文件
    with open(kb_path, 'w', encoding='utf-8') as f:
        json.dump(kb_list, f, ensure_ascii=False, indent=2)
    
    action_text = "更新" if found else "添加"
    return {"code": 0, "msg": f"已成功{action_text}到知识库"}

@app.post("/auth/updateUserInfo")
async def update_user_info(data: UpdateUserInfoRequest, payload: dict = Depends(verify_token)):
    async with SessionLocal() as session:
        result = await session.execute(select(User).where(User.id == int(payload["sub"])))  # sub转int
        user = result.scalar_one_or_none()
        if not user:
            raise HTTPException(404, "用户不存在")
        user.real_name = data.realName
        user.nick_name = data.nickName
        await session.commit()
        return {"code": 0, "msg": "信息已更新"}

@app.post("/auth/updatePreferences")
async def update_preferences(data: UpdatePreferencesRequest, payload: dict = Depends(verify_token)):
    async with SessionLocal() as session:
        result = await session.execute(select(User).where(User.id == int(payload["sub"])))  # sub转int
        user = result.scalar_one_or_none()
        if not user:
            raise HTTPException(404, "用户不存在")
        user.theme = data.theme
        user.language = data.language
        await session.commit()
        return {"code": 0, "msg": "偏好设置已保存"}

# 创建会议接口
@app.post("/api/meetings/create")
async def create_meeting(data: MeetingCreateRequest, payload: dict = Depends(verify_token)):
    user_id = payload["userId"]
    
    # 查找主持人用户ID（如果host是用户名）
    host_user_id = 0
    if data.host:
        async with SessionLocal() as session:
            result = await session.execute(select(User).where(User.username == data.host))
            host_user = result.scalar_one_or_none()
            if host_user:
                host_user_id = host_user.id
    
    # 处理参会人列表
    participants_str = ",".join(map(str, data.participants)) if data.participants else ""
    
    async with MeetingSessionLocal() as session:
        new_meeting = Meeting(
            title=data.title,
            host=data.host,
            time=data.time,
            location=data.location,
            period=data.period,
            status=data.status,
            user_id=user_id,
            participants=participants_str,
            host_user_id=host_user_id or user_id  # 如果找不到主持人，默认为创建者
        )
        session.add(new_meeting)
        await session.commit()
        return {"code": 0, "msg": "会议已安排", "meetingId": new_meeting.id}

# 删除会议接口
@app.post("/api/meetings/delete")
async def delete_meeting(data: dict = Body(...), payload: dict = Depends(verify_token)):
    meeting_id = data.get("id")
    if not meeting_id:
        return {"code": 1, "msg": "未提供会议ID"}
    
    user_id = payload["userId"]
    try:
        async with MeetingSessionLocal() as session:
            result = await session.execute(select(Meeting).where(Meeting.id == int(meeting_id)))
            meeting = result.scalar_one_or_none()
            
            if not meeting:
                return {"code": 1, "msg": "会议不存在"}
            
            # 只有创建者或主持人才可以删除会议
            if meeting.user_id != user_id and meeting.host_user_id != user_id:
                return {"code": 1, "msg": "无权限删除该会议"}
            
            await session.delete(meeting)
            await session.commit()
            return {"code": 0, "msg": "会议删除成功"}
    except Exception as e:
        print(f"删除会议失败: {e}")
        return {"code": 1, "msg": "删除会议失败"}

# 会议列表接口
@app.get("/api/meetings/list")
async def list_meetings(
    page: int = 1, 
    pageSize: int = 10, 
    status: str = None,
    kw: str = None,
    payload: dict = Depends(verify_token)
):
    try:
        user_id = payload["userId"]
        async with MeetingSessionLocal() as session:
            # 查询用户相关的会议：创建者、主持人、参会人
            try:
                result = await session.execute(select(Meeting))
                all_meetings = result.scalars().all()
            except Exception as e:
                print(f"[list_meetings] 查询会议失败: {e}")
                return {"code": 1, "msg": "获取会议列表失败", "data": {"list": [], "total": 0}}
            
            # 过滤出用户相关的会议
            user_meetings = []
            for meeting in all_meetings:
                try:
                    is_related = False
                    
                    # 1. 是创建者
                    if meeting.user_id and meeting.user_id == user_id:
                        is_related = True
                    
                    # 2. 是主持人
                    elif meeting.host_user_id and meeting.host_user_id == user_id:
                        is_related = True
                    
                    # 3. 是参会人
                    elif meeting.participants:
                        try:
                            participant_ids = [int(pid.strip()) for pid in meeting.participants.split(",") if pid.strip().isdigit()]
                            if user_id in participant_ids:
                                is_related = True
                        except Exception:
                            pass  # 忽略解析错误
                    
                    if is_related:
                        user_meetings.append(meeting)
                except Exception as e:
                    print(f"[获取会议列表] 处理会议 {meeting.id} 时出错: {e}")
                    continue
            
            # 根据状态过滤（如果提供了状态参数且不为空）
            if status and status.strip():
                user_meetings = [m for m in user_meetings if (m.status or "upcoming") == status]
            
            # 根据关键词搜索（如果提供了关键词且不为空）
            if kw and kw.strip():
                kw_lower = kw.strip().lower()
                user_meetings = [
                    m for m in user_meetings 
                    if (m.title and kw_lower in m.title.lower()) or 
                       (m.host and kw_lower in m.host.lower())
                ]
            
            # 分页处理
            total = len(user_meetings)
            start = (page - 1) * pageSize
            end = start + pageSize
            data = []
            
            for m in user_meetings[start:end]:
                try:
                    # 解析参会人ID列表
                    participant_ids = []
                    if m.participants:
                        try:
                            participant_ids = [int(pid.strip()) for pid in m.participants.split(",") if pid.strip().isdigit()]
                        except Exception:
                            participant_ids = []
                    
                    data.append({
                        "id": m.id,
                        "title": m.title or "",
                        "host": m.host or "",
                        "time": m.time or "",
                        "location": m.location or "",
                        "period": m.period or "",
                        "status": m.status or "upcoming",
                        "participants": participant_ids,
                        "host_user_id": getattr(m, 'host_user_id', 0) or 0
                    })
                except Exception as e:
                    print(f"[获取会议列表] 序列化会议 {m.id} 时出错: {e}")
                    continue
            
            return {"code": 0, "data": {"list": data, "total": total}}
    except Exception as e:
        print(f"[获取会议列表] 发生错误: {e}")
        import traceback
        traceback.print_exc()
        return {"code": 1, "msg": f"获取会议列表失败: {str(e)}", "data": {"list": [], "total": 0}}

# 会议室列表接口（静态返回）
@app.get("/api/meetings/rooms")
async def get_rooms():
    rooms = [
        {"id": 1, "name": "会议室A"},
        {"id": 2, "name": "会议室B"},
        {"id": 3, "name": "会议室C"},
        {"id": 4, "name": "3131"},
    ]
    return {"code": 0, "data": rooms}

# 会议冲突检测接口
@app.post("/api/meetings/check-conflict")
async def check_conflict(data: dict = Body(...)):
    time = data.get("time")
    location = data.get("location")
    async with MeetingSessionLocal() as session:
        result = await session.execute(select(Meeting).where(Meeting.time == time, Meeting.location == location))
        conflict = result.scalar_one_or_none() is not None
        return {"code": 0, "data": {"conflict": conflict}}

# 编辑会议接口
@app.post("/api/meetings/edit")
async def edit_meeting(data: dict = Body(...), payload: dict = Depends(verify_token)):
    user_id = payload["userId"]
    meeting_id = data.get("id")
    if not meeting_id:
        return {"code": 1, "msg": "缺少会议ID"}
    
    async with MeetingSessionLocal() as session:
        result = await session.execute(select(Meeting).where(Meeting.id == meeting_id))
        meeting = result.scalar_one_or_none()
        if not meeting:
            return {"code": 1, "msg": "会议不存在"}
        
        # 检查权限：创建者、主持人或参会人才能编辑
        has_permission = False
        
        # 1. 是创建者
        if meeting.user_id == user_id:
            has_permission = True
        
        # 2. 是主持人
        elif meeting.host_user_id == user_id:
            has_permission = True
        
        # 3. 是参会人
        elif meeting.participants:
            participant_ids = [int(pid.strip()) for pid in meeting.participants.split(",") if pid.strip().isdigit()]
            if user_id in participant_ids:
                has_permission = True
        
        if not has_permission:
            return {"code": 1, "msg": "无权限编辑此会议"}
        
        # 更新字段
        meeting.title = data.get("title", meeting.title)
        meeting.host = data.get("host", meeting.host)
        meeting.time = data.get("time", meeting.time)
        meeting.location = data.get("location", meeting.location)
        meeting.period = data.get("period", meeting.period)
        meeting.status = data.get("status", meeting.status)
        
        # 处理参会人更新
        if "participants" in data:
            participants = data.get("participants", [])
            if isinstance(participants, list):
                meeting.participants = ",".join(map(str, participants))
            else:
                meeting.participants = str(participants)
        
        # 更新主持人用户ID
        if data.get("host"):
            async with SessionLocal() as user_session:
                result = await user_session.execute(select(User).where(User.username == data.get("host")))
                host_user = result.scalar_one_or_none()
                if host_user:
                    meeting.host_user_id = host_user.id
        
        await session.commit()
        return {"code": 0, "msg": "保存成功"}

@app.post("/api/assistant/analyze-file")
async def analyze_file(data: dict = Body(...), payload: dict = Depends(verify_token)):
    """分析上传的文件内容"""
    file_url = data.get("file_url", "")
    user_question = data.get("question", "")
    
    if not file_url:
        return {"code": 1, "msg": "文件URL不能为空"}
    
    # 从环境变量读取API KEY，如果没有则使用默认值
    QWEN_API_KEY = os.getenv("QWEN_API_KEY", "sk-751689471abb4aaf9d7169c86884b1a0")
    
    try:
        # 解析文件路径
        if file_url.startswith('/uploads/assistant/'):
            uploads_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'uploads')
            file_path = os.path.join(uploads_dir, 'assistant', file_url.split('assistant/')[-1])
        else:
            return {"code": 1, "msg": "不支持的文件路径格式"}
        
        if not os.path.exists(file_path):
            return {"code": 1, "msg": "文件不存在"}
        
        # 提取文件内容
        file_ext = os.path.splitext(file_path)[1]
        file_content = extract_file_content(file_path, file_ext)
        
        if not file_content or file_content.startswith("[错误") or file_content.startswith("[不支持"):
            return {"code": 1, "msg": f"无法提取文件内容: {file_content}"}
        
        # 限制内容长度（避免超过API限制）
        max_content_length = 30000  # 约30000字符
        if len(file_content) > max_content_length:
            file_content = file_content[:max_content_length] + "\n\n[内容已截断，仅显示前30000字符]"
        
        # 构建提示词
        if not user_question or user_question.strip() == "":
            prompt = f"""请详细分析以下文件内容，包括：
1. 文件的主要内容和主题
2. 关键信息和要点
3. 文件的结构和组织方式
4. 任何重要的数据、数字或统计信息
5. 总结和建议

文件内容：
{file_content}

请用简体中文回答，提供详细、结构化的分析。"""
        else:
            prompt = f"""请根据以下问题分析文件内容：

问题：{user_question}

文件内容：
{file_content}

请用简体中文回答，提供详细、准确的分析。"""
        
        # 调用通义千问API
        url = "https://dashscope.aliyuncs.com/api/v1/services/aigc/text-generation/generation"
        headers = {
            "Authorization": f"Bearer {QWEN_API_KEY}",
            "Content-Type": "application/json"
        }
        
        payload = {
            "model": "qwen-turbo",
            "input": {
                "prompt": prompt
            },
            "parameters": {
                "result_format": "message"
            }
        }
        
        resp = requests.post(url, headers=headers, json=payload, timeout=60)
        
        # 检查HTTP状态码
        if resp.status_code != 200:
            return {"code": 1, "msg": f"API调用失败，状态码: {resp.status_code}"}
        
        result = resp.json()
        
        # 检查是否有错误
        if "error" in result:
            return {"code": 1, "msg": f"API错误: {result.get('error', {}).get('message', '未知错误')}"}
        
        # 解析返回结果
        reply = result.get("output", {}).get("text")
        if not reply:
            choices = result.get("output", {}).get("choices")
            if choices and isinstance(choices, list) and len(choices) > 0:
                reply = choices[0].get("message", {}).get("content")
        
        # 确保reply是字符串
        if reply:
            reply_str = str(reply) if reply else "分析完成，但未返回内容"
            return {"code": 0, "data": {"reply": reply_str, "content_length": len(file_content)}}
        else:
            error_msg = result.get("message", "文件分析失败，API未返回有效内容")
            return {"code": 1, "msg": str(error_msg), "debug": result}
    
    except Exception as e:
        import traceback
        error_detail = traceback.format_exc()
        return {"code": 1, "msg": f"文件分析异常: {str(e)}", "detail": error_detail}

@app.post("/api/assistant/chat")
async def assistant_chat(data: dict = Body(...)):
    user_message = data.get("message", "")
    image_url = data.get("image_url", "")  # 获取图片URL
    # 从环境变量读取API KEY，如果没有则使用默认值
    QWEN_API_KEY = os.getenv("QWEN_API_KEY", "sk-751689471abb4aaf9d7169c86884b1a0")
    
    # 如果有图片，使用多模态模型
    if image_url:
        # 使用通义千问多模态模型 qwen-vl-max
        url = "https://dashscope.aliyuncs.com/api/v1/services/aigc/multimodal-generation/generation"
        headers = {
            "Authorization": f"Bearer {QWEN_API_KEY}",
            "Content-Type": "application/json"
        }
        
        # 构建完整的图片URL（如果是相对路径，需要转换为绝对路径）
        if image_url.startswith('/uploads/'):
            # 如果是相对路径，需要根据实际情况构建完整URL
            # 这里假设图片在服务器上可以直接访问
            full_image_url = image_url
        elif image_url.startswith('http'):
            full_image_url = image_url
        else:
            full_image_url = f"/uploads/{image_url}"
        
        # 如果没有提供问题，使用默认提示
        if not user_message or user_message.strip() == "":
            user_message = "请详细分析这张图片，包括图片中的内容、文字、场景等所有可见信息。"
        
        # 构建多模态消息
        prompt = f"请用简体中文回答：{user_message}"
        
        # 读取图片并转换为base64
        try:
            image_base64 = None
            # 尝试读取本地图片文件并转换为base64
            if full_image_url.startswith('/uploads/'):
                uploads_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'uploads')
                # 处理assistant子目录的情况
                if 'assistant' in full_image_url:
                    image_path = os.path.join(uploads_dir, 'assistant', full_image_url.split('assistant/')[-1])
                elif 'images' in full_image_url:
                    image_path = os.path.join(uploads_dir, 'images', full_image_url.split('images/')[-1])
                else:
                    image_path = os.path.join(uploads_dir, full_image_url.replace('/uploads/', ''))
                
                if os.path.exists(image_path):
                    with open(image_path, 'rb') as f:
                        image_data = f.read()
                        image_base64 = base64.b64encode(image_data).decode('utf-8')
                else:
                    return {"code": 1, "msg": f"图片文件不存在: {image_path}"}
            elif full_image_url.startswith('http'):
                # 如果是HTTP URL，尝试下载并转换为base64
                try:
                    resp = requests.get(full_image_url, timeout=10)
                    if resp.status_code == 200:
                        image_base64 = base64.b64encode(resp.content).decode('utf-8')
                    else:
                        return {"code": 1, "msg": f"无法下载图片: HTTP {resp.status_code}"}
                except Exception as e:
                    return {"code": 1, "msg": f"下载图片失败: {str(e)}"}
            else:
                return {"code": 1, "msg": "不支持的图片URL格式"}
            
            if not image_base64:
                return {"code": 1, "msg": "无法获取图片数据"}
            
            # 使用通义千问多模态API (qwen-vl-plus 或 qwen-vl-max)
            # 通义千问VL API使用messages格式，支持图片base64
            payload = {
                "model": "qwen-vl-plus",  # 使用qwen-vl-plus，如果API KEY支持qwen-vl-max可以改用
                "input": {
                    "messages": [
                        {
                            "role": "user",
                            "content": [
                                {
                                    "image": f"data:image/jpeg;base64,{image_base64}"
                                },
                                {
                                    "text": prompt
                                }
                            ]
                        }
                    ]
                },
                "parameters": {
                    "result_format": "message"
                }
            }
            
            resp = requests.post(url, headers=headers, json=payload, timeout=60)
            result = resp.json()
            
            # 解析返回结果 - 通义千问VL的返回格式
            if "output" in result:
                choices = result.get("output", {}).get("choices", [])
                if choices and len(choices) > 0:
                    message_content = choices[0].get("message", {}).get("content", "")
                    if message_content:
                        reply = message_content
                    else:
                        # 尝试其他格式
                        reply = result.get("output", {}).get("text", "")
                else:
                    reply = result.get("output", {}).get("text", "")
            else:
                # 如果返回格式不同，尝试直接获取
                reply = result.get("text", "") or result.get("message", {}).get("content", "")
            
            if reply:
                # 确保reply是字符串
                reply_str = str(reply) if reply else "图片分析完成，但未返回内容"
                return {"code": 0, "data": {"reply": reply_str}}
            else:
                # 检查是否是账户问题
                error_msg = result.get("message", "图片分析失败")
                if "Access denied" in str(result) or "overdue" in str(result).lower() or "account" in str(result).lower():
                    error_msg = "API账户异常，请检查账户状态和余额。如需更换API KEY，请设置环境变量 QWEN_API_KEY"
                return {"code": 1, "msg": error_msg, "debug": result}
                
        except Exception as e:
            # 如果多模态API失败，返回详细错误信息
            import traceback
            error_detail = traceback.format_exc()
            return {"code": 1, "msg": f"图片分析异常: {str(e)}", "detail": error_detail}
    else:
        # 纯文本模式，使用原来的逻辑
        url = "https://dashscope.aliyuncs.com/api/v1/services/aigc/text-generation/generation"
        headers = {
            "Authorization": f"Bearer {QWEN_API_KEY}",
            "Content-Type": "application/json"
        }
        # 强制中文回答
        prompt = f"请用简体中文回答：{user_message}"
        payload = {
            "model": "qwen-turbo",
            "input": {
                "prompt": prompt
            },
            "parameters": {
                "result_format": "message"
            }
        }
        try:
            resp = requests.post(url, headers=headers, json=payload, timeout=30)
            result = resp.json()
            # 通义千问返回格式
            reply = result.get("output", {}).get("text")
            if not reply:
                # 兼容新版API返回格式
                choices = result.get("output", {}).get("choices")
                if choices and isinstance(choices, list):
                    reply = choices[0].get("message", {}).get("content")
            if reply:
                return {"code": 0, "data": {"reply": reply}}
            else:
                # 检查是否是账户问题
                error_msg = result.get("message", "通义千问API调用失败")
                if "Access denied" in str(result) or "overdue" in str(result).lower() or "account" in str(result).lower():
                    error_msg = "API账户异常，请检查账户状态和余额。如需更换API KEY，请设置环境变量 QWEN_API_KEY"
                return {"code": 1, "msg": error_msg, "debug": result}
        except Exception as e:
            error_msg = str(e)
            if "Access denied" in error_msg or "overdue" in error_msg.lower():
                error_msg = "API账户异常，请检查账户状态和余额。如需更换API KEY，请设置环境变量 QWEN_API_KEY"
            return {"code": 1, "msg": f"通义千问API异常: {error_msg}"}

@app.get("/api/assistant/recommend")
async def assistant_recommend(query: str = ""):
    # 这里可以根据query做智能推荐，先返回模拟数据
    return [
        {"id": 1, "title": "如何高效开会", "type": "知识"},
        {"id": 2, "title": "本周工作总结模板", "type": "模板"},
        {"id": 3, "title": "会议纪要写作技巧", "type": "技巧"}
    ]

@app.get("/api/schedule/today")
async def schedule_today(payload: dict = Depends(verify_token)):
    """获取今天的日程 - 只返回当前用户的数据"""
    try:
        user_id = payload["userId"]
        today_str = datetime.now().strftime('%Y-%m-%d')
        schedule_data = []
        
        # 获取今天的任务
        async with SessionLocal() as session:
            result = await session.execute(
                select(TodayTask).where(
                    TodayTask.user_id == user_id,
                    TodayTask.date == today_str
                )
            )
            tasks = result.scalars().all()
            for task in tasks:
                schedule_data.append({
                    "id": f"task_{task.id}",
                    "title": task.content,
                    "time": task.time or "待定",
                    "type": "task"
                })
        
        # 获取今天的会议
        async with MeetingSessionLocal() as session:
            result = await session.execute(
                select(Meeting).where(Meeting.time.like(f"{today_str}%"))
            )
            all_meetings = result.scalars().all()
            
            # 过滤出用户相关的会议
            for meeting in all_meetings:
                is_related = False
                if meeting.user_id == user_id or meeting.host_user_id == user_id:
                    is_related = True
                elif meeting.participants:
                    participant_ids = [int(pid.strip()) for pid in meeting.participants.split(",") if pid.strip().isdigit()]
                    if user_id in participant_ids:
                        is_related = True
                
                if is_related:
                    # 提取时间部分
                    meeting_time = meeting.time.split(' ')[1] if ' ' in meeting.time else meeting.time
                    schedule_data.append({
                        "id": f"meeting_{meeting.id}",
                        "title": meeting.title,
                        "time": meeting_time,
                        "type": "meeting"
                    })
        
        # 按时间排序
        schedule_data.sort(key=lambda x: x.get("time", "00:00"))
        
        return {
            "code": 0,
            "data": schedule_data
        }
    except Exception as e:
        print(f"[schedule_today] 错误: {e}")
        return {"code": 1, "msg": "获取日程失败", "data": []}

@app.post("/api/user/upload-avatar")
async def upload_avatar(file: UploadFile = File(...), payload: dict = Depends(verify_token)):
    user_id = payload["userId"]
    ext = os.path.splitext(file.filename)[-1]
    save_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'uploads', 'avatar')
    os.makedirs(save_dir, exist_ok=True)
    save_path = os.path.join(save_dir, f"{user_id}{ext}")
    with open(save_path, "wb") as f:
        f.write(await file.read())
    avatar_url = f"/uploads/avatar/{user_id}{ext}"
    async with SessionLocal() as session:
        result = await session.execute(select(User).where(User.id == user_id))
        user = result.scalar_one_or_none()
        if user:
            user.avatar = avatar_url
            await session.commit()
    return {"code": 0, "msg": "上传成功", "data": {"avatar": avatar_url}}

@app.post("/api/user/generate-avatar")
async def generate_new_avatar(payload: dict = Depends(verify_token)):
    """为当前用户生成新的随机头像"""
    user_id = payload["userId"]
    
    try:
        async with SessionLocal() as session:
            result = await session.execute(select(User).where(User.id == user_id))
            user = result.scalar_one_or_none()
            if not user:
                return {"code": 1, "msg": "用户不存在"}
            
            # 生成新的随机头像
            avatar_url = generate_random_avatar()
            
            # 如果头像是外部URL，尝试下载并保存到本地
            if avatar_url.startswith('https://'):
                try:
                    local_avatar = save_avatar_from_url(avatar_url, user_id)
                    user.avatar = local_avatar
                except Exception as e:
                    print(f"下载头像失败，使用默认头像: {e}")
                    user.avatar = "avatar/default.svg"
            else:
                user.avatar = avatar_url
            
            await session.commit()
            
            return {
                "code": 0, 
                "msg": "头像生成成功", 
                "data": {"avatar": user.avatar}
            }
    except Exception as e:
        return {"code": 1, "msg": f"生成头像失败: {str(e)}"}

@app.post("/api/upload/image")
async def upload_image(file: UploadFile = File(...), payload: dict = Depends(verify_token)):
    """上传图片接口，用于聊天发送图片"""
    if file.content_type not in ["image/jpeg", "image/png", "image/gif", "image/webp", "image/svg+xml"]:
        raise HTTPException(status_code=400, detail="不支持的文件类型，仅支持 JPEG、PNG、GIF、WebP、SVG 格式")
    
    user_id = payload.get("userId")
    if not user_id:
        raise HTTPException(status_code=400, detail="用户ID不能为空")
    
    # 文件大小限制 5MB
    if file.size and file.size > 5 * 1024 * 1024:
        raise HTTPException(status_code=400, detail="文件大小不能超过 5MB")
    
    # 创建目录
    uploads_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'uploads', 'images')
    os.makedirs(uploads_dir, exist_ok=True)
    
    # 生成文件名
    timestamp = int(time.time())
    file_ext = os.path.splitext(file.filename)[1] if file.filename else ".jpg"
    filename = f"chat_{user_id}_{timestamp}{file_ext}"
    file_path = os.path.join(uploads_dir, filename)
    
    # 保存文件
    with open(file_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
    
    return {"code": 0, "data": {"url": f"/uploads/images/{filename}"}}

def extract_file_content(file_path: str, file_ext: str) -> str:
    """提取文件内容，支持多种文件格式"""
    content = ""
    try:
        ext = file_ext.lower()
        
        if ext == ".txt":
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read()
        
        elif ext == ".docx":
            try:
                from docx import Document
                doc = Document(file_path)
                content = "\n".join([p.text for p in doc.paragraphs])
            except ImportError:
                content = "[错误: 需要安装python-docx库]"
            except Exception as e:
                content = f"[解析DOCX文件时出错: {str(e)}]"
        
        elif ext == ".doc":
            # 处理旧版Word格式(.doc)
            try:
                # 尝试使用textract库（如果已安装）
                import textract
                content = textract.process(file_path).decode('utf-8')
            except ImportError:
                # 如果textract未安装，尝试使用pypandoc
                try:
                    import pypandoc
                    content = pypandoc.convert_file(file_path, 'plain', format='doc')
                except ImportError:
                    # 如果都未安装，提示用户安装
                    content = "[错误: 需要安装textract或pypandoc库来解析.doc文件。请运行: pip install textract 或 pip install pypandoc]"
                except Exception as e:
                    content = f"[解析DOC文件时出错: {str(e)}]"
            except Exception as e:
                content = f"[解析DOC文件时出错: {str(e)}]"
        
        elif ext == ".pdf":
            try:
                import PyPDF2
                with open(file_path, 'rb') as f:
                    pdf_reader = PyPDF2.PdfReader(f)
                    pages_text = []
                    for page_num in range(len(pdf_reader.pages)):
                        page = pdf_reader.pages[page_num]
                        pages_text.append(page.extract_text())
                    content = "\n\n".join(pages_text)
            except ImportError:
                content = "[错误: 需要安装PyPDF2库]"
            except Exception as e:
                content = f"[解析PDF文件时出错: {str(e)}]"
        
        elif ext == ".pptx":
            try:
                from pptx import Presentation
                prs = Presentation(file_path)
                slides = []
                for slide in prs.slides:
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            slide_text.append(shape.text)
                    if slide_text:
                        slides.append("\n".join(slide_text))
                content = "\n\n---幻灯片分隔---\n\n".join(slides)
            except ImportError:
                content = "[错误: 需要安装python-pptx库]"
            except Exception as e:
                content = f"[解析PPTX文件时出错: {str(e)}]"
        
        elif ext in [".xls", ".xlsx"]:
            try:
                import openpyxl
                wb = openpyxl.load_workbook(file_path)
                sheets_content = []
                for sheet_name in wb.sheetnames:
                    sheet = wb[sheet_name]
                    rows = []
                    for row in sheet.iter_rows(values_only=True):
                        row_data = [str(cell) if cell is not None else "" for cell in row]
                        rows.append("\t".join(row_data))
                    sheets_content.append(f"工作表: {sheet_name}\n" + "\n".join(rows))
                content = "\n\n---工作表分隔---\n\n".join(sheets_content)
            except ImportError:
                content = "[错误: 需要安装openpyxl库]"
            except Exception as e:
                content = f"[解析Excel文件时出错: {str(e)}]"
        
        elif ext == ".csv":
            try:
                import csv
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    reader = csv.reader(f)
                    rows = []
                    for row in reader:
                        rows.append("\t".join(row))
                    content = "\n".join(rows)
            except Exception as e:
                content = f"[解析CSV文件时出错: {str(e)}]"
        
        else:
            content = f"[不支持的文件格式: {ext}]"
    
    except Exception as e:
        content = f"[提取文件内容时出错: {str(e)}]"
    
    return content

@app.post("/api/assistant/upload")
async def assistant_upload_file(file: UploadFile = File(...), payload: dict = Depends(verify_token)):
    """智能助手上传接口，支持图片和附件"""
    
    user_id = payload.get("userId")
    if not user_id:
        raise HTTPException(status_code=400, detail="用户ID不能为空")
    
    # 文件大小限制 10MB
    if file.size and file.size > 10 * 1024 * 1024:
        raise HTTPException(status_code=400, detail="文件大小不能超过 10MB")
    
    # 检查文件类型（根据文件扩展名和content_type）
    file_ext = os.path.splitext(file.filename)[1].lower() if file.filename else ""
    file_type = "file"
    
    # 优先根据文件扩展名判断
    if file_ext in [".jpg", ".jpeg", ".png", ".gif", ".webp", ".svg", ".bmp"]:
        file_type = "image"
    elif file_ext in [".txt", ".md", ".log"]:
        file_type = "text"
    elif file_ext in [".doc", ".docx", ".pdf", ".ppt", ".pptx"]:
        file_type = "document"
    elif file_ext in [".xls", ".xlsx", ".csv"]:
        file_type = "document"  # Excel文件也归类为document，可以分析
    elif file_ext in [".zip", ".rar", ".tar", ".gz", ".7z"]:
        file_type = "archive"
    elif file.content_type:
        # 如果扩展名无法判断，使用content_type
        if file.content_type.startswith("image/"):
            file_type = "image"
        elif file.content_type in ["application/pdf", "application/msword", "application/vnd.openxmlformats-officedocument.wordprocessingml.document"]:
            file_type = "document"
        elif file.content_type in ["application/vnd.ms-excel", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", "text/csv"]:
            file_type = "document"  # Excel文件
        elif file.content_type in ["application/zip", "application/x-rar-compressed", "application/x-tar", "application/gzip"]:
            file_type = "archive"
        elif file.content_type.startswith("text/"):
            file_type = "text"
    
    # 创建目录
    uploads_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'uploads', 'assistant')
    os.makedirs(uploads_dir, exist_ok=True)
    
    # 生成文件名
    timestamp = int(time.time())
    file_ext = os.path.splitext(file.filename)[1] if file.filename else ""
    filename = f"assistant_{user_id}_{timestamp}{file_ext}"
    file_path = os.path.join(uploads_dir, filename)
    
    # 保存文件
    with open(file_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
    
    # 如果是文档类型，尝试提取内容
    content = ""
    if file_type in ["document", "text"]:
        content = extract_file_content(file_path, file_ext)
    
    return {
        "code": 0, 
        "data": {
            "url": f"/uploads/assistant/{filename}",
            "type": file_type,
            "original_name": file.filename,
            "size": file.size or 0,
            "content": content[:5000] if content else ""  # 返回前5000字符，用于预览
        }
    }

class TodayTaskModel(BaseModel):
    content: str
    time: str = ""
    completed: bool = False
    date: str = ""
    type: str = "bg-primary"
    endDate: str = ""  # 新增，前端传 endDate

@app.get("/api/today-tasks")
async def get_today_tasks(payload: dict = Depends(verify_token)):
    try:
        user_id = payload["userId"]
        async with SessionLocal() as session:
            # 只查询当前用户的任务，处理 user_id 可能为 NULL 的情况
            try:
                result = await session.execute(
                    select(TodayTask).where(
                        (TodayTask.user_id == user_id) | (TodayTask.user_id.is_(None))
                    )
                )
                tasks = result.scalars().all()
                # 过滤掉 user_id 为 NULL 的任务（这些是旧数据，应该被迁移）
                tasks = [t for t in tasks if t.user_id == user_id]
                tasks = [dict(id=t.id, content=t.content, time=t.time, completed=t.completed, date=t.date, type=getattr(t, "type", "bg-primary"), endDate=getattr(t, "end_date", "")) for t in tasks]
            except Exception as db_error:
                # 如果查询失败，可能是字段不存在，尝试使用原始 SQL
                print(f"[获取今日任务] ORM 查询失败，尝试使用原始 SQL 查询: {db_error}")
                result = await session.execute(
                    text("SELECT * FROM today_tasks WHERE user_id = :user_id"),
                    {"user_id": user_id}
                )
                rows = result.fetchall()
                tasks = []
                for row in rows:
                    tasks.append({
                        "id": row[0],
                        "content": row[2] if len(row) > 2 else "",
                        "time": row[3] if len(row) > 3 else "待定",
                        "completed": bool(row[4]) if len(row) > 4 else False,
                        "date": row[5] if len(row) > 5 else "",
                        "type": row[6] if len(row) > 6 else "bg-primary",
                        "endDate": row[7] if len(row) > 7 else ""
                    })
        return {"code": 0, "tasks": tasks}
    except Exception as e:
        print(f"[get_today_tasks] 错误: {e}")
        import traceback
        traceback.print_exc()
        return {"code": 1, "msg": f"获取任务失败: {str(e)}", "tasks": []}

@app.post("/api/today-tasks")
async def add_today_task(task: TodayTaskModel, payload: dict = Depends(verify_token)):
    try:
        user_id = payload["userId"]
        import datetime
        db_path = os.path.abspath(os.path.join(os.path.dirname(__file__), '../users.db'))
        print(f"[add_today_task] 数据库绝对路径: {db_path}")
        print(f"[add_today_task] 收到任务: {task}, 用户ID: {user_id}")
        async with SessionLocal() as session:
            new_task = TodayTask(
                user_id=user_id,  # 添加用户ID
                content=task.content,
                time=task.time or datetime.datetime.now().strftime('%H:%M'),
                completed=task.completed,
                date=task.date or datetime.datetime.now().strftime('%Y-%m-%d'),
                type=task.type or "bg-primary",
                end_date=task.endDate or ""
            )
            session.add(new_task)
            await session.commit()
            await session.refresh(new_task)
            print(f"[add_today_task] 已写入数据库: {new_task}")
            return {"code": 0, "msg": "添加成功", "task": dict(id=new_task.id, content=new_task.content, time=new_task.time, completed=new_task.completed, date=new_task.date, type=new_task.type, endDate=new_task.end_date)}
    except Exception as e:
        print(f"[add_today_task] 错误: {e}")
        return {"code": 1, "msg": "添加任务失败"}

@app.delete("/api/today-tasks/{task_id}")
async def delete_today_task(task_id: int, payload: dict = Depends(verify_token)):
    try:
        user_id = payload["userId"]
        print(f"[delete_today_task] 请求删除任务ID: {task_id}, 用户ID: {user_id}")
        async with SessionLocal() as session:
            # 只允许删除自己的任务
            result = await session.execute(select(TodayTask).where(TodayTask.id == task_id, TodayTask.user_id == user_id))
            task = result.scalar_one_or_none()
            if not task:
                print(f"[delete_today_task] 未找到任务或无权删除: {task_id}")
                return JSONResponse({"code": 1, "msg": "任务不存在或无权限"}, status_code=status.HTTP_404_NOT_FOUND)
            await session.delete(task)
            await session.commit()
            print(f"[delete_today_task] 已删除任务: {task_id}")
            return {"code": 0, "msg": "删除成功"}
    except Exception as e:
        print(f"[delete_today_task] 错误: {e}")
        return {"code": 1, "msg": "删除任务失败"}

@app.put("/api/today-tasks/{task_id}")
async def update_today_task(task_id: int = Path(...), task: TodayTaskModel = Body(...), payload: dict = Depends(verify_token)):
    try:
        user_id = payload["userId"]
        async with SessionLocal() as session:
            # 只允许更新自己的任务
            result = await session.execute(select(TodayTask).where(TodayTask.id == task_id, TodayTask.user_id == user_id))
            db_task = result.scalar_one_or_none()
            if not db_task:
                return JSONResponse({"code": 1, "msg": "任务不存在或无权限"}, status_code=status.HTTP_404_NOT_FOUND)
            db_task.content = task.content
            db_task.time = task.time
            db_task.completed = task.completed
            db_task.date = task.date
            db_task.type = task.type or "bg-primary"
            db_task.end_date = task.endDate or ""
            await session.commit()
            await session.refresh(db_task)
            return {"code": 0, "msg": "更新成功", "task": {
                "id": db_task.id,
                "content": db_task.content,
                "time": db_task.time,
                "completed": db_task.completed,
                "date": db_task.date,
                "type": db_task.type,
                "endDate": db_task.end_date
            }}
    except Exception as e:
        print(f"[update_today_task] 错误: {e}")
        return {"code": 1, "msg": "更新任务失败"}

app.mount("/uploads", StaticFiles(directory=os.path.join(os.path.dirname(os.path.dirname(__file__)), 'uploads')), name="uploads")

@app.post("/api/ai/summary")
async def ai_summary(data: dict = Body(...)):
    content = data.get("content", "")
    summary_type = data.get("type", "general")  # general, meeting_minutes
    meeting_title = data.get("meetingTitle", "")
    
    if not content:
        return {"code": 1, "msg": "内容不能为空"}
    
    # 根据类型构建不同的提示词
    if summary_type == "meeting_minutes":
        prompt = f"""请根据以下会议录音转录内容，生成一份专业的会议纪要：

会议主题：{meeting_title or '未指定'}
录音转录内容：
{content}

请按以下格式生成会议纪要：
1. 会议基本信息
2. 会议要点
3. 决议事项
4. 待办事项
5. 下次会议安排（如有）

要求：
- 提取关键信息和决策点
- 语言简洁明了
- 结构化呈现
- 重点突出行动项"""
    else:
        prompt = f"请对以下内容进行简明扼要的中文摘要：\n{content}"
    
    # 尝试使用通义千问API
    try:
        # 从环境变量读取API KEY，如果没有则使用默认值
        api_key = os.getenv("QWEN_API_KEY", "sk-751689471abb4aaf9d7169c86884b1a0")
        api_url = "https://dashscope.aliyuncs.com/api/v1/services/aigc/text-generation/generation"
        headers = {
            "Authorization": f"Bearer {api_key}",
            "Content-Type": "application/json"
        }
        payload = {
            "model": "qwen-turbo",
            "input": {"prompt": prompt},
            "parameters": {"result_format": "message"}
        }
        async with httpx.AsyncClient() as client:
            resp = await client.post(api_url, headers=headers, json=payload, timeout=30)
            result = resp.json()
            # 解析通义千问返回的摘要内容
            summary = result.get("output", {}).get("text", "")
            if summary:
                print("通义千问返回：", result)
                return {"code": 0, "summary": summary}
    except Exception as e:
        print(f"通义千问API调用失败: {e}")
    
    # 如果API调用失败，使用本地摘要算法
    try:
        if summary_type == "meeting_minutes":
            summary = generate_meeting_minutes(content, meeting_title)
        else:
            summary = generate_local_summary(content)
        return {"code": 0, "summary": summary}
    except Exception as e:
        print(f"本地摘要生成失败: {e}")
        return {"code": 1, "msg": "摘要生成失败"}

def generate_local_summary(content: str) -> str:
    """
    本地智能摘要生成算法
    """
    import re
    
    # 移除HTML标签
    clean_content = re.sub(r'<[^>]+>', '', content)
    
    # 分句
    sentences = re.split(r'[。！？；\n]', clean_content)
    sentences = [s.strip() for s in sentences if s.strip() and len(s.strip()) > 5]
    
    if not sentences:
        return "文档内容过短，无法生成摘要。"
    
    # 关键词提取和重要句子识别
    keywords = ['项目', '系统', '功能', '技术', '应用', '创新', '管理', '服务', '用户', '数据', 
                '智能', '平台', '解决方案', '优化', '提升', '实现', '支持', '开发', '设计', '架构']
    
    # 句子评分
    sentence_scores = []
    for sentence in sentences:
        score = 0
        # 关键词得分
        for keyword in keywords:
            if keyword in sentence:
                score += 1
        # 位置得分（开头和结尾的句子得分更高）
        if sentences.index(sentence) < len(sentences) * 0.3:
            score += 2
        elif sentences.index(sentence) > len(sentences) * 0.7:
            score += 1
        # 长度得分
        if 10 <= len(sentence) <= 50:
            score += 1
        
        sentence_scores.append((sentence, score))
    
    # 排序并选择重要句子
    sentence_scores.sort(key=lambda x: x[1], reverse=True)
    
    # 选择前3-5个重要句子
    top_sentences = sentence_scores[:min(5, len(sentence_scores))]
    
    # 生成摘要
    summary_parts = []
    
    # 添加标题或主题
    if '标题' in clean_content or 'h1' in content.lower():
        title_match = re.search(r'<h1[^>]*>(.*?)</h1>', content, re.IGNORECASE)
        if title_match:
            summary_parts.append(f"📋 主题：{title_match.group(1)}")
    
    # 添加核心内容
    summary_parts.append("📝 核心内容：")
    for sentence, score in top_sentences:
        if sentence and len(sentence) > 8:
            summary_parts.append(f"• {sentence}")
            if len(summary_parts) >= 8:  # 控制摘要长度
                break
    
    # 添加统计信息
    word_count = len(clean_content)
    summary_parts.append(f"\n📊 文档统计：约{word_count}字，{len(sentences)}句话")
    
    # 如果内容包含特定结构，添加结构化信息
    if '<table' in content:
        summary_parts.append("📋 包含表格数据")
    if '<ul' in content or '<ol' in content:
        summary_parts.append("📋 包含列表信息")
    if '<h2' in content or '<h3' in content:
        summary_parts.append("📋 包含章节结构")
    
    summary = "\n".join(summary_parts)
    
    # 如果摘要太短，添加更多信息
    if len(summary) < 100:
        summary += f"\n\n💡 这是一个关于{extract_main_topic(clean_content)}的文档。"
    
    return summary

def extract_main_topic(content: str) -> str:
    """提取文档主题"""
    topics = {
        '项目': ['项目', '系统', '平台', '软件', '应用'],
        '管理': ['管理', '计划', '安排', '组织', '协调'],
        '技术': ['技术', '开发', '编程', '代码', '算法'],
        '教育': ['教育', '学习', '培训', '课程', '教学'],
        '政策': ['政策', '规定', '制度', '法规', '文件'],
        '报告': ['报告', '总结', '分析', '统计', '调研']
    }
    
    for topic, keywords in topics.items():
        if any(keyword in content for keyword in keywords):
            return topic
    
    return "综合内容"

def generate_meeting_minutes(content: str, meeting_title: str = "") -> str:
    """
    生成会议纪要
    """
    import re
    from datetime import datetime
    
    # 移除HTML标签
    clean_content = re.sub(r'<[^>]+>', '', content)
    
    # 分句
    sentences = re.split(r'[。！？；\n]', clean_content)
    sentences = [s.strip() for s in sentences if s.strip() and len(s.strip()) > 5]
    
    if not sentences:
        return "录音内容过短，无法生成有效纪要。"
    
    # 关键词分类
    keywords = {
        '决策': ['决定', '确定', '同意', '批准', '通过', '决议'],
        '计划': ['计划', '安排', '准备', '预定', '打算'],
        '问题': ['问题', '困难', '挑战', '风险', '障碍'],
        '建议': ['建议', '提议', '推荐', '意见', '看法'],
        '行动': ['执行', '实施', '完成', '负责', '跟进'],
        '时间': ['今天', '明天', '下周', '下个月', '截止', '期限'],
        '人员': ['负责人', '团队', '部门', '同事', '参与者']
    }
    
    # 分类句子
    categorized = {category: [] for category in keywords.keys()}
    other_content = []
    
    for sentence in sentences:
        classified = False
        for category, words in keywords.items():
            if any(word in sentence for word in words):
                categorized[category].append(sentence)
                classified = True
                break
        if not classified:
            other_content.append(sentence)
    
    # 生成纪要
    minutes_parts = []
    
    # 头部信息
    minutes_parts.append(f"<h2>会议纪要</h2>")
    minutes_parts.append(f"<p><strong>会议主题：</strong>{meeting_title or '未指定'}</p>")
    minutes_parts.append(f"<p><strong>纪要生成时间：</strong>{datetime.now().strftime('%Y年%m月%d日 %H:%M')}</p>")
    minutes_parts.append("<hr>")
    
    # 主要内容
    if categorized['决策']:
        minutes_parts.append("<h3>📋 会议决议</h3>")
        for i, item in enumerate(categorized['决策'][:5], 1):
            minutes_parts.append(f"<p>{i}. {item}</p>")
    
    if categorized['计划']:
        minutes_parts.append("<h3>📅 工作计划</h3>")
        for i, item in enumerate(categorized['计划'][:5], 1):
            minutes_parts.append(f"<p>{i}. {item}</p>")
    
    if categorized['问题']:
        minutes_parts.append("<h3>⚠️ 讨论问题</h3>")
        for i, item in enumerate(categorized['问题'][:3], 1):
            minutes_parts.append(f"<p>{i}. {item}</p>")
    
    if categorized['建议']:
        minutes_parts.append("<h3>💡 建议意见</h3>")
        for i, item in enumerate(categorized['建议'][:3], 1):
            minutes_parts.append(f"<p>{i}. {item}</p>")
    
    if categorized['行动']:
        minutes_parts.append("<h3>✅ 行动事项</h3>")
        for i, item in enumerate(categorized['行动'][:5], 1):
            minutes_parts.append(f"<p>{i}. {item}</p>")
    
    # 其他重要内容
    if other_content:
        minutes_parts.append("<h3>📝 其他内容</h3>")
        important_others = [s for s in other_content if len(s) > 15][:3]
        for i, item in enumerate(important_others, 1):
            minutes_parts.append(f"<p>{i}. {item}</p>")
    
    # 如果没有足够的结构化内容，提供基本摘要
    if len([item for items in categorized.values() for item in items]) < 3:
        minutes_parts.append("<h3>📄 会议摘要</h3>")
        # 选择最重要的句子
        important_sentences = sorted(sentences, key=len, reverse=True)[:5]
        for i, sentence in enumerate(important_sentences, 1):
            minutes_parts.append(f"<p>{i}. {sentence}</p>")
    
    # 添加统计信息
    minutes_parts.append("<hr>")
    minutes_parts.append(f"<p><small>📊 会议时长：约{len(clean_content)//200}分钟 | 转录字数：{len(clean_content)}字</small></p>")
    
    return "\n".join(minutes_parts)

# 新增会议附件上传接口
ATTACHMENTS_PATH = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'uploads', 'attachments.json')

def load_attachments():
    if not os.path.exists(ATTACHMENTS_PATH):
        return []
    with open(ATTACHMENTS_PATH, 'r', encoding='utf-8') as f:
        return json.load(f)

def save_attachments(data):
    with open(ATTACHMENTS_PATH, 'w', encoding='utf-8') as f:
        json.dump(data, f, ensure_ascii=False, indent=2)

@app.post("/api/meetings/upload")
async def upload_meeting_attachment(meetingId: int = Form(...), file: UploadFile = File(...), payload: dict = Depends(verify_token)):
    uploads_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'uploads', 'meetings')
    os.makedirs(uploads_dir, exist_ok=True)
    save_path = os.path.join(uploads_dir, f"{meetingId}_{file.filename}")
    with open(save_path, "wb") as f:
        f.write(await file.read())
    url = f"/uploads/meetings/{meetingId}_{file.filename}"
    # 记录到 attachments.json
    attachments = load_attachments()
    attachments.append({
        "meetingId": meetingId,
        "filename": file.filename,
        "url": url,
        "upload_time": datetime.now().strftime('%Y-%m-%d %H:%M:%S')
    })
    save_attachments(attachments)
    return {"code": 0, "msg": "上传成功", "filename": file.filename, "meetingId": meetingId, "url": url}

@app.get("/api/meetings/{meeting_id}/attachments")
async def get_meeting_attachments(meeting_id: int):
    attachments = load_attachments()
    return [a for a in attachments if int(a["meetingId"]) == int(meeting_id)]

@app.get("/api/meetings/approval-history")
async def approval_history(id: int):
    # 目前没有审批历史，返回空列表即可
    return []

@app.post("/api/reset-password")
async def reset_password(data: ResetPasswordRequest):
    # 校验验证码
    real_code = r.get(f"captcha:{data.contact}")
    if not real_code:
        return {"code": 1, "msg": "验证码已过期或不存在"}
    if data.captcha != real_code:
        return {"code": 2, "msg": "验证码错误"}
    r.delete(f"captcha:{data.contact}")

    # 修改密码
    async with SessionLocal() as session:
        result = await session.execute(select(User).where(User.contact == data.contact))
        user = result.scalar_one_or_none()
        if not user:
            return {"code": 3, "msg": "用户不存在"}
        import bcrypt
        user.password = bcrypt.hashpw(data.new_password.encode(), bcrypt.gensalt()).decode()
        await session.commit()
        return {"code": 0, "msg": "密码重置成功"}

# ========== WebSocket 单聊消息推送与存储 =============
import json as _json
active_connections = {}

@app.websocket("/ws/chat/user")
async def websocket_endpoint(websocket: WebSocket, token: str):
    user_id = None
    try:
        # 接受所有来源的连接（支持跨域）
        await websocket.accept()
        
        payload = jwt.decode(token, SECRET_KEY, algorithms=[ALGORITHM])
        user_id = payload.get("userId") or payload.get("sub")
        if not user_id:
            print(f"[WS] token解析失败，payload={payload}")
            await websocket.close(code=1008, reason="Invalid token")
            return
        user_id = int(user_id)
        print(f"[WS] 用户 {user_id} 连接成功，来源: {websocket.client}")
        
        # 如果该用户已有连接，关闭旧连接
        if user_id in active_connections:
            try:
                await active_connections[user_id].close(code=1000, reason="New connection")
            except Exception as e:
                print(f"[WS] 关闭旧连接异常: {e}")
        
        active_connections[user_id] = websocket
        print(f"[WS] 当前活跃连接数: {len(active_connections)}")
        while True:
            try:
                data = await websocket.receive_text()
                try:
                    msg_obj = _json.loads(data)
                except Exception as e:
                    print(f"[WS] JSON解析失败: {e}, data={data}")
                    continue
                # 存储消息到 chat-db.json
                try:
                    save_message_to_json(msg_obj)
                except Exception as e:
                    print(f"[WS] 消息存储异常: {e}, msg={msg_obj}")
                # 推送给目标用户和自己
                to_user_id = msg_obj.get("to_user_id")
                # 推送给自己
                if user_id in active_connections:
                    try:
                        await active_connections[user_id].send_text(_json.dumps(msg_obj))
                    except Exception as e:
                        print(f"[WS] 推送给自己异常: {e}")
                # 推送给对方（如果不是自己）
                if to_user_id in active_connections and to_user_id != user_id:
                    try:
                        await active_connections[to_user_id].send_text(_json.dumps(msg_obj))
                    except Exception as e:
                        print(f"[WS] 推送给对方异常: {e}")
            except WebSocketDisconnect:
                print(f"[WS] 用户 {user_id} 断开连接")
                break
            except Exception as e:
                print(f"[WS] 消息处理异常: {e}")
    finally:
        if user_id and user_id in active_connections:
            del active_connections[user_id]


@app.post("/api/message/send")
async def send_message(data: dict = Body(...), payload: dict = Depends(verify_token)):
    """发送消息API接口（非WebSocket方式）"""
    from_user_id = payload["userId"]
    to_user_id = data.get("to_user_id")
    content = data.get("content")
    msg_type = data.get("type", "text")
    file_name = data.get("file_name")  # 文件名
    file_size = data.get("file_size")  # 文件大小
    
    if not to_user_id or not content:
        return {"code": 1, "msg": "缺少必要参数"}
    
    msg_obj = {
        "from_user_id": from_user_id,
        "to_user_id": to_user_id,
        "content": content,
        "type": msg_type,
        "time": datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    }
    
    # 如果是文件类型，添加文件元数据
    if msg_type == "file":
        if file_name:
            msg_obj["file_name"] = file_name
        if file_size is not None:
            msg_obj["file_size"] = file_size
    
    try:
        save_message_to_json(msg_obj)
        return {"code": 0, "msg": "发送成功"}
    except Exception as e:
        return {"code": 1, "msg": f"发送失败: {str(e)}"}

@app.post("/api/message/read")
async def mark_message_read(data: dict = Body(...), payload: dict = Depends(verify_token)):
    """标记消息为已读"""
    msg_id = data.get("msg_id")
    if not msg_id:
        return {"code": 1, "msg": "缺少消息ID"}
    
    # 从chat-db.json读取并更新消息状态
    db_path = os.path.join(os.path.dirname(__file__), "chat-db.json")
    
    if os.path.exists(db_path):
        with open(db_path, "r", encoding="utf-8") as f:
            db = _json.load(f)
            
        # 查找并更新消息
        for room_key, messages in db.get("rooms", {}).items():
            for msg in messages:
                if msg.get("id") == msg_id:
                    msg["is_read"] = 1
                    
                    # 保存更新后的数据
                    with open(db_path, "w", encoding="utf-8") as f:
                        _json.dump(db, f, ensure_ascii=False, indent=2)
                    
                    return {"code": 0, "msg": "标记成功"}
    
    return {"code": 1, "msg": "消息不存在"}

@app.get("/api/message/unread_counts")
async def get_unread_counts(payload: dict = Depends(verify_token)):
    """获取个人聊天未读消息计数"""
    user_id = payload["userId"]
    
    # 从chat-db.json读取所有消息
    db_path = os.path.join(os.path.dirname(__file__), "chat-db.json")
    unread_counts = {}
    
    if os.path.exists(db_path):
        with open(db_path, "r", encoding="utf-8") as f:
            db = _json.load(f)
            
        # 统计每个发送者的未读消息数
        for room_key, messages in db.get("rooms", {}).items():
            for msg in messages:
                # 如果是发给当前用户的消息，且未读（没有is_read字段或is_read为0）
                # 支持字符串和数字类型的用户ID比较
                to_user_id = msg.get("to_user_id")
                if (to_user_id == user_id or str(to_user_id) == str(user_id)) and msg.get("is_read", 0) == 0:
                    from_user_id = msg.get("from_user_id")
                    if from_user_id:
                        # 确保键为字符串类型
                        key = str(from_user_id)
                        unread_counts[key] = unread_counts.get(key, 0) + 1
    
    print(f"用户 {user_id} 的未读消息计数: {unread_counts}")  # 调试信息
    return {"code": 0, "data": unread_counts}

@app.get("/api/group/unread_counts")
async def get_group_unread_counts(payload: dict = Depends(verify_token)):
    """获取群聊未读消息计数"""
    user_id = payload["userId"]
    
    # 这里应该从群聊消息数据库读取，暂时返回空
    # 实际实现需要根据群聊消息的存储方式来统计
    unread_counts = {}
    
    return {"code": 0, "data": unread_counts}

def save_message_to_json(msg_obj):
    import os
    db_path = os.path.join(os.path.dirname(__file__), "chat-db.json")
    if not isinstance(msg_obj, dict):
        raise ValueError("消息对象必须是字典类型")
    for k in ["from_user_id", "to_user_id", "content", "time"]:
        if k not in msg_obj:
            raise ValueError(f"消息对象缺少必要字段: {k}")
    
    # 为消息添加id和is_read字段
    msg_obj["id"] = int(datetime.now().timestamp() * 1000)  # 使用时间戳作为ID
    msg_obj["is_read"] = 0  # 默认未读
    
    if os.path.exists(db_path):
        with open(db_path, "r", encoding="utf-8") as f:
            db = _json.load(f)
    else:
        db = {"rooms": {}}
    # 只在 user_{from_user_id} 存储消息，避免重复
    room_key = f"user_{msg_obj['from_user_id']}"
    db["rooms"].setdefault(room_key, []).append(msg_obj)
    with open(db_path, "w", encoding="utf-8") as f:
        _json.dump(db, f, ensure_ascii=False, indent=2)
    
    # 删除自动创建消息通知的功能

# 删除了 create_message_notification 函数

# 删除了 get_user_info_by_id 函数

# 删除所有通知相关API

@app.get("/api/users/list")
async def get_users_list(payload: dict = Depends(verify_token)):
    """获取用户列表，用于选择参会人"""
    async with SessionLocal() as session:
        result = await session.execute(select(User))
        users = result.scalars().all()
        
        user_list = []
        for user in users:
            user_list.append({
                "id": user.id,
                "username": user.username,
                "realName": user.real_name or user.username,
                "nickName": user.nick_name or user.username,
                "department": user.department,
                "avatar": user.avatar,
                "contact": user.contact
            })
        
        return {"code": 0, "data": user_list}

# ========== 人事部人员管理接口 ==========
@app.get("/api/personnel/list")
async def get_personnel_list(
    page: int = 1,
    pageSize: int = 20,
    department: str = None,
    keyword: str = None,
    payload: dict = Depends(verify_token)
):
    """获取人员列表（支持分页、部门筛选、关键词搜索）"""
    try:
        current_user_id = payload["userId"]
        # 检查当前用户是否是人事部
        async with SessionLocal() as session:
            result = await session.execute(select(User).where(User.id == current_user_id))
            current_user = result.scalar_one_or_none()
            if not current_user or current_user.department != "人事部":
                return {"code": 1, "msg": "无权限访问，仅人事部可查看", "data": {"list": [], "total": 0}}
            
            # 构建查询
            query = select(User)
            
            # 部门筛选
            if department and department.strip():
                query = query.where(User.department == department.strip())
            
            # 关键词搜索（姓名、用户名、联系方式）
            if keyword and keyword.strip():
                keyword_lower = keyword.strip().lower()
                query = query.where(
                    or_(
                        User.real_name.like(f"%{keyword_lower}%"),
                        User.username.like(f"%{keyword_lower}%"),
                        User.contact.like(f"%{keyword_lower}%"),
                        User.nick_name.like(f"%{keyword_lower}%")
                    )
                )
            
            # 执行查询
            result = await session.execute(query)
            all_users = result.scalars().all()
            
            # 分页处理
            total = len(all_users)
            start = (page - 1) * pageSize
            end = start + pageSize
            paginated_users = all_users[start:end]
            
            # 格式化返回数据
            user_list = []
            for user in paginated_users:
                user_list.append({
                    "id": user.id,
                    "username": user.username,
                    "realName": user.real_name or "",
                    "nickName": user.nick_name or "",
                    "department": user.department or "",
                    "avatar": user.avatar or "",
                    "contact": user.contact or "",
                    "theme": user.theme or "auto",
                    "language": user.language or "zh"
                })
            
            return {"code": 0, "data": {"list": user_list, "total": total}}
    except Exception as e:
        print(f"[get_personnel_list] 错误: {e}")
        return {"code": 1, "msg": f"获取人员列表失败: {str(e)}", "data": {"list": [], "total": 0}}

@app.post("/api/personnel/create")
async def create_personnel(data: dict = Body(...), payload: dict = Depends(verify_token)):
    """创建新人员"""
    try:
        current_user_id = payload["userId"]
        async with SessionLocal() as session:
            # 检查权限
            result = await session.execute(select(User).where(User.id == current_user_id))
            current_user = result.scalar_one_or_none()
            if not current_user or current_user.department != "人事部":
                return {"code": 1, "msg": "无权限，仅人事部可创建人员"}
            
            # 验证必填字段
            username = data.get("username", "").strip()
            contact = data.get("contact", "").strip()
            password = data.get("password", "").strip()
            
            if not username:
                return {"code": 1, "msg": "用户名不能为空"}
            if not contact:
                return {"code": 1, "msg": "联系方式不能为空"}
            if not password:
                return {"code": 1, "msg": "密码不能为空"}
            
            # 检查用户名是否已存在
            result = await session.execute(select(User).where(User.username == username))
            if result.scalar_one_or_none():
                return {"code": 1, "msg": "用户名已存在"}
            
            # 检查联系方式是否已存在
            result = await session.execute(select(User).where(User.contact == contact))
            if result.scalar_one_or_none():
                return {"code": 1, "msg": "联系方式已存在"}
            
            # 创建新用户
            hashed_password = bcrypt.hashpw(password.encode(), bcrypt.gensalt()).decode()
            new_user = User(
                username=username,
                password=hashed_password,
                contact=contact,
                real_name=data.get("realName", "").strip(),
                nick_name=data.get("nickName", "").strip(),
                department=data.get("department", "").strip(),
                avatar=data.get("avatar", ""),
                theme=data.get("theme", "auto"),
                language=data.get("language", "zh")
            )
            session.add(new_user)
            await session.commit()
            await session.refresh(new_user)
            
            return {"code": 0, "msg": "创建成功", "data": {"id": new_user.id}}
    except Exception as e:
        print(f"[create_personnel] 错误: {e}")
        return {"code": 1, "msg": f"创建失败: {str(e)}"}

@app.post("/api/personnel/update")
async def update_personnel(data: dict = Body(...), payload: dict = Depends(verify_token)):
    """更新人员信息"""
    try:
        current_user_id = payload["userId"]
        user_id = data.get("id")
        if not user_id:
            return {"code": 1, "msg": "用户ID不能为空"}
        
        async with SessionLocal() as session:
            # 检查权限
            result = await session.execute(select(User).where(User.id == current_user_id))
            current_user = result.scalar_one_or_none()
            if not current_user or current_user.department != "人事部":
                return {"code": 1, "msg": "无权限，仅人事部可更新人员信息"}
            
            # 查找要更新的用户
            result = await session.execute(select(User).where(User.id == user_id))
            user = result.scalar_one_or_none()
            if not user:
                return {"code": 1, "msg": "用户不存在"}
            
            # 更新字段
            if "username" in data:
                new_username = data["username"].strip()
                if new_username and new_username != user.username:
                    # 检查新用户名是否已存在
                    result = await session.execute(select(User).where(User.username == new_username))
                    if result.scalar_one_or_none():
                        return {"code": 1, "msg": "用户名已存在"}
                    user.username = new_username
            
            if "contact" in data:
                new_contact = data["contact"].strip()
                if new_contact and new_contact != user.contact:
                    # 检查新联系方式是否已存在
                    result = await session.execute(select(User).where(User.contact == new_contact))
                    if result.scalar_one_or_none():
                        return {"code": 1, "msg": "联系方式已存在"}
                    user.contact = new_contact
            
            if "realName" in data:
                user.real_name = data["realName"].strip()
            if "nickName" in data:
                user.nick_name = data["nickName"].strip()
            if "department" in data:
                user.department = data["department"].strip()
            if "avatar" in data:
                user.avatar = data["avatar"]
            if "theme" in data:
                user.theme = data["theme"]
            if "language" in data:
                user.language = data["language"]
            
            # 如果提供了新密码，更新密码
            if "password" in data and data["password"]:
                hashed_password = bcrypt.hashpw(data["password"].encode(), bcrypt.gensalt()).decode()
                user.password = hashed_password
            
            await session.commit()
            return {"code": 0, "msg": "更新成功"}
    except Exception as e:
        print(f"[update_personnel] 错误: {e}")
        return {"code": 1, "msg": f"更新失败: {str(e)}"}

@app.post("/api/personnel/delete")
async def delete_personnel(data: dict = Body(...), payload: dict = Depends(verify_token)):
    """删除人员"""
    try:
        current_user_id = payload["userId"]
        user_id = data.get("id")
        if not user_id:
            return {"code": 1, "msg": "用户ID不能为空"}
        
        if int(user_id) == int(current_user_id):
            return {"code": 1, "msg": "不能删除自己"}
        
        async with SessionLocal() as session:
            # 检查权限
            result = await session.execute(select(User).where(User.id == current_user_id))
            current_user = result.scalar_one_or_none()
            if not current_user or current_user.department != "人事部":
                return {"code": 1, "msg": "无权限，仅人事部可删除人员"}
            
            # 查找要删除的用户
            result = await session.execute(select(User).where(User.id == user_id))
            user = result.scalar_one_or_none()
            if not user:
                return {"code": 1, "msg": "用户不存在"}
            
            await session.delete(user)
            await session.commit()
            return {"code": 0, "msg": "删除成功"}
    except Exception as e:
        print(f"[delete_personnel] 错误: {e}")
        return {"code": 1, "msg": f"删除失败: {str(e)}"}

@app.get("/api/dashboard/stats")
async def get_dashboard_stats(payload: dict = Depends(verify_token)):
    """获取工作台统计数据 - 只统计当前用户的数据"""
    try:
        user_id = payload["userId"]
        
        # 获取文档统计 - 只统计当前用户的文档
        doc_path = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'uploads', 'doc.json')
        doc_total = 0
        if os.path.exists(doc_path):
            try:
                with open(doc_path, 'r', encoding='utf-8') as f:
                    docs = json.load(f)
                    # 只统计当前用户的文档
                    doc_total = len([doc for doc in docs if doc.get("userId") == user_id])
            except:
                doc_total = 0
        
        # 获取任务统计 - 只统计当前用户的任务
        async with SessionLocal() as session:
            result = await session.execute(select(func.count(TodayTask.id)).where(TodayTask.user_id == user_id))
            task_total = result.scalar() or 0
            
            # 获取今天的会议数量 - 只统计当前用户相关的会议
            today_str = datetime.now().strftime('%Y-%m-%d')
            async with MeetingSessionLocal() as meeting_session:
                # 查询所有今天的会议
                result = await meeting_session.execute(
                    select(Meeting).where(Meeting.time.like(f"{today_str}%"))
            )
                all_today_meetings = result.scalars().all()
                
                # 过滤出用户相关的会议
                user_today_meetings = []
                for meeting in all_today_meetings:
                    is_related = False
                    if meeting.user_id == user_id or meeting.host_user_id == user_id:
                        is_related = True
                    elif meeting.participants:
                        participant_ids = [int(pid.strip()) for pid in meeting.participants.split(",") if pid.strip().isdigit()]
                        if user_id in participant_ids:
                            is_related = True
                    if is_related:
                        user_today_meetings.append(meeting)
                
                today_meetings = len(user_today_meetings)
            
            # 获取即将开始的会议（今天且未来时间）
            current_time = datetime.now().strftime('%Y-%m-%d %H:%M')
            upcoming_meetings = 0
            completed_meetings = 0
            for meeting in user_today_meetings:
                if meeting.time >= current_time:
                    upcoming_meetings += 1
                else:
                    completed_meetings += 1
        
        # 生成基于真实数据的统计 - 只统计当前用户的数据
        base_visits = 1000
        base_views = 800
        
        stats = {
            "cardStats": [
                {
                    "des": "总访问次数",
                    "icon": "&#xe721;",
                    "num": base_visits + (doc_total * 25) + (task_total * 10),
                    "change": "+15%"
                },
                {
                    "des": "我的任务",
                    "icon": "&#xe724;",
                    "num": task_total,
                    "change": "+8%"
                },
                {
                    "des": "点击量",
                    "icon": "&#xe7aa;",
                    "num": base_views + (doc_total * 40) + (today_meetings * 25),
                    "change": "-5%"
                },
                {
                    "des": "我的文档",
                    "icon": "&#xe82a;",
                    "num": doc_total,
                    "change": "+25%"
                }
            ],
            "docStats": {
                "total": doc_total,
                "todayProcessed": min(max(doc_total // 30, 5), 60),
                "pending": min(max(doc_total // 60, 2), 15)
            },
            "scheduleStats": {
                "todaySchedule": today_meetings,
                "upcoming": upcoming_meetings,
                "completed": completed_meetings
            }
        }
        
        return {"code": 0, "data": stats}
        
    except Exception as e:
        print(f"获取工作台统计数据失败: {e}")
        # 返回默认数据
        return {"code": 0, "data": {
            "cardStats": [
                {"des": "总访问次数", "icon": "&#xe721;", "num": 5000, "change": "+15%"},
                {"des": "在线访客数", "icon": "&#xe724;", "num": 50, "change": "+8%"},
                {"des": "点击量", "icon": "&#xe7aa;", "num": 4500, "change": "-5%"},
                {"des": "新用户", "icon": "&#xe82a;", "num": 30, "change": "+25%"}
            ],
            "docStats": {"total": 0, "todayProcessed": 0, "pending": 0},
            "scheduleStats": {"todaySchedule": 0, "upcoming": 0, "completed": 0}
        }}

@app.get("/api/dashboard/chart-data")
async def get_dashboard_chart_data(payload: dict = Depends(verify_token)):
    """获取工作台图表数据"""
    try:
        # 生成过去9个月的业务数据
        chart_data = []
        months = ['一月', '二月', '三月', '四月', '五月', '六月', '七月', '八月', '九月']
        
        # 基于当前数据生成模拟的月度数据
        async with SessionLocal() as session:
            result = await session.execute(select(func.count(Meeting.id)))
            meeting_count = result.scalar() or 0
            
            result = await session.execute(select(func.count(User.id)))
            user_count = result.scalar() or 0
            
        # 计算基础值
        base_value = max(meeting_count * 8 + user_count * 3, 60)
        
        # 生成有趋势的月度数据
        for i, month in enumerate(months):
            # 模拟业务增长趋势
            trend_factor = 1 + (i * 0.1)  # 每月增长10%
            seasonal_factor = 1 + (0.3 * (i % 3) / 3)  # 季节性波动
            random_factor = 1 + (random.randint(-20, 30) / 100)  # 随机波动
            
            value = int(base_value * trend_factor * seasonal_factor * random_factor)
            chart_data.append({
                "month": month,
                "value": max(value, 30)  # 确保最小值为30
            })
        
        return {"code": 0, "data": chart_data}
        
    except Exception as e:
        print(f"获取图表数据失败: {e}")
        # 返回默认数据
        return {"code": 0, "data": [
            {"month": "一月", "value": 80},
            {"month": "二月", "value": 120},
            {"month": "三月", "value": 150},
            {"month": "四月", "value": 100},
            {"month": "五月", "value": 90},
            {"month": "六月", "value": 130},
            {"month": "七月", "value": 140},
            {"month": "八月", "value": 170},
            {"month": "九月", "value": 160}
        ]}

@app.get("/api/dashboard/department")
async def get_department_info(payload: dict = Depends(verify_token)):
    """获取部门信息"""
    try:
        user_id = payload.get('sub')
        
        # 获取当前用户信息
        async with SessionLocal() as session:
            result = await session.execute(select(User).where(User.id == user_id))
            user = result.scalar_one_or_none()
            
            if not user:
                return {"code": 404, "msg": "用户不存在"}
            
            # 获取同部门的其他用户
            result = await session.execute(
                select(User).where(
                    User.department == user.department,
                    User.id != user_id
                )
            )
            department_users = result.scalars().all()
            
        # 根据部门生成相应的部门信息
        department_configs = {
            "人事部": {
                "responsibilities": [
                    "负责公司招聘、员工入职与离职管理，制定招聘计划，筛选简历，组织面试",
                    "组织员工培训与绩效考核，建立培训体系，制定KPI指标，定期评估员工表现",
                    "管理员工档案与薪酬福利，维护人事档案，制定薪酬方案，处理社保公积金",
                    "推动企业文化建设与员工关系维护，组织团建活动，处理员工投诉，营造良好工作氛围"
                ],
                "email": "hr@company.com",
                "phone": "010-12345678",
                "announcement": "本月员工生日会将于25日举行，请大家准时参加！"
            },
            "技术部": {
                "responsibilities": [
                    "负责系统开发与技术架构设计，制定技术方案，编写核心代码，确保系统稳定性",
                    "维护服务器运行与数据安全，监控系统性能，备份重要数据，防范安全风险",
                    "推动技术创新与产品优化，研究新技术应用，优化用户体验，提升产品竞争力",
                    "提供技术支持与培训服务，解决技术问题，培训其他部门员工，编写技术文档"
                ],
                "email": "tech@company.com",
                "phone": "010-12345679",
                "announcement": "新版系统将于下周上线，请各位做好测试准备！"
            },
            "市场部": {
                "responsibilities": [
                    "制定市场营销策略与推广方案，分析市场趋势，制定营销计划，执行推广活动",
                    "管理客户关系与商务合作，维护重要客户，拓展新客户，建立合作伙伴关系",
                    "分析市场趋势与竞争对手，收集市场信息，分析竞品动态，提供决策建议",
                    "组织品牌宣传与活动策划，提升品牌知名度，策划营销活动，管理品牌形象"
                ],
                "email": "marketing@company.com",
                "phone": "010-12345680",
                "announcement": "Q4季度市场推广活动即将启动，欢迎大家踊跃参与！"
            },
            "财务部": {
                "responsibilities": [
                    "负责公司财务核算与报表编制，处理日常账务，编制财务报表，确保财务数据准确",
                    "管理资金流动与成本控制，制定资金计划，控制成本支出，优化资金使用效率",
                    "处理税务申报与合规事务，按时申报纳税，处理税务问题，确保合规经营",
                    "提供财务分析与决策支持，分析财务数据，提供决策建议，支持业务发展"
                ],
                "email": "finance@company.com",
                "phone": "010-12345681",
                "announcement": "月度财务报表提交截止日期为本月30日，请各部门配合！"
            },
            "运营部": {
                "responsibilities": [
                    "负责产品运营与用户增长，制定运营策略，提升用户活跃度，实现用户增长目标",
                    "管理内容运营与社区建设，策划内容活动，维护用户社区，提升用户粘性",
                    "分析运营数据与效果评估，监控关键指标，分析运营效果，优化运营策略",
                    "协调各部门资源与项目推进，统筹项目进度，协调资源分配，确保项目顺利执行"
                ],
                "email": "operations@company.com",
                "phone": "010-12345682",
                "announcement": "本月运营数据表现良好，用户活跃度提升15%，继续保持！"
            },
            "销售部": {
                "responsibilities": [
                    "制定销售策略与目标管理，制定销售计划，设定销售目标，跟踪销售进度",
                    "开发新客户与维护老客户，拓展销售渠道，维护客户关系，提升客户满意度",
                    "管理销售团队与业绩考核，培训销售技能，考核销售业绩，激励团队士气",
                    "分析销售数据与市场反馈，收集客户反馈，分析销售趋势，优化销售策略"
                ],
                "email": "sales@company.com",
                "phone": "010-12345683",
                "announcement": "本月销售目标完成率95%，距离目标还有5%，大家加油！"
            },
            "客服部": {
                "responsibilities": [
                    "提供客户服务与问题解决，处理客户咨询，解决客户问题，提升客户满意度",
                    "管理客户投诉与纠纷处理，及时响应投诉，妥善处理纠纷，维护公司形象",
                    "收集客户反馈与需求分析，整理客户建议，分析客户需求，为产品改进提供依据",
                    "建立客户档案与关系维护，建立客户数据库，定期回访客户，维护长期合作关系"
                ],
                "email": "service@company.com",
                "phone": "010-12345684",
                "announcement": "客户满意度调查即将开始，请各位同事积极配合！"
            },
            "法务部": {
                "responsibilities": [
                    "处理法律事务与合同审查，审查业务合同，处理法律纠纷，提供法律咨询",
                    "管理知识产权与商标注册，保护公司知识产权，处理商标专利事务，防范法律风险",
                    "制定合规制度与风险控制，建立合规体系，识别法律风险，制定防范措施",
                    "提供法律培训与合规指导，培训员工法律知识，指导业务合规，提升法律意识"
                ],
                "email": "legal@company.com",
                "phone": "010-12345685",
                "announcement": "新版合同模板已更新，请各部门使用最新版本！"
            }
        }
        
        # 获取当前部门配置
        dept_config = department_configs.get(user.department, department_configs["人事部"])
        
        # 确定部门领导和成员
        leader = user.real_name or user.username
        members = []
        
        for dept_user in department_users:
            display_name = dept_user.real_name or dept_user.username
            members.append(display_name)
        
        # 如果有多个用户，第一个作为领导，其他作为成员
        if len(members) > 0:
            if len(members) == 1:
                # 如果只有一个其他用户，当前用户作为领导
                pass
            else:
                # 如果有多个用户，第一个作为领导，当前用户和其他作为成员
                leader = members[0]
                members = [user.real_name or user.username] + members[1:]
        
        department_info = {
            "name": user.department,
            "responsibilities": dept_config["responsibilities"],
            "leader": leader,
            "members": members,
            "email": dept_config["email"],
            "phone": dept_config["phone"],
            "announcement": dept_config["announcement"]
        }
        
        return {"code": 0, "data": department_info}
        
    except Exception as e:
        print(f"获取部门信息失败: {e}")
        # 返回默认数据
        return {"code": 0, "data": {
            "name": "人事部",
            "responsibilities": [
                "负责公司招聘、员工入职与离职管理",
                "组织员工培训与绩效考核",
                "管理员工档案与薪酬福利",
                "推动企业文化建设与员工关系维护"
            ],
            "leader": "张三",
            "members": ["李四", "王五"],
            "email": "hr@company.com",
            "phone": "010-12345678",
            "announcement": "本月员工生日会将于25日举行，请大家准时参加！"
        }}

# ========== 录音历史相关接口 ==========
@app.post("/api/recordings/upload")
async def upload_recording(file: UploadFile = File(...), data: str = Form(...), payload: dict = Depends(verify_token)):
    """上传录音文件并保存记录"""
    user_id = payload["userId"]
    
    try:
        # 解析表单数据
        import json
        recording_data = json.loads(data)
        
        # 创建录音文件存储目录
        recordings_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'uploads', 'recordings')
        os.makedirs(recordings_dir, exist_ok=True)
        
        # 生成唯一文件名
        import uuid
        file_extension = os.path.splitext(file.filename)[1] if file.filename else '.webm'
        unique_filename = f"recording_{user_id}_{int(time.time())}_{str(uuid.uuid4())[:8]}{file_extension}"
        file_path = os.path.join(recordings_dir, unique_filename)
        
        # 保存文件
        with open(file_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)
        
        # 获取文件大小
        file_size = os.path.getsize(file_path)
        
        # 保存录音记录到数据库
        async with SessionLocal() as session:
            new_recording = RecordingHistory(
                meeting_id=recording_data.get("meeting_id"),
                meeting_title=recording_data.get("meeting_title", ""),
                user_id=user_id,
                filename=file.filename or unique_filename,
                file_path=f"recordings/{unique_filename}",
                transcript=recording_data.get("transcript", ""),
                minutes=recording_data.get("minutes", ""),
                duration=recording_data.get("duration", 0),
                file_size=file_size,
                created_at=datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
                status="completed"
            )
            session.add(new_recording)
            await session.commit()
            await session.refresh(new_recording)
            
            return {
                "code": 0, 
                "msg": "录音上传成功", 
                "data": {
                    "id": new_recording.id,
                    "filename": new_recording.filename,
                    "file_path": new_recording.file_path
                }
            }
    
    except Exception as e:
        print(f"录音上传失败: {e}")
        return {"code": 1, "msg": f"录音上传失败: {str(e)}"}

@app.get("/api/recordings/history")
async def get_recording_history(page: int = 1, pageSize: int = 10, keyword: str = "", payload: dict = Depends(verify_token)):
    """获取录音历史列表"""
    user_id = payload["userId"]
    
    try:
        async with SessionLocal() as session:
            # 构建查询
            query = select(RecordingHistory).where(RecordingHistory.user_id == user_id)
            
            # 关键词搜索
            if keyword:
                query = query.where(
                    RecordingHistory.meeting_title.contains(keyword) |
                    RecordingHistory.filename.contains(keyword)
                )
            
            # 按创建时间倒序
            query = query.order_by(RecordingHistory.id.desc())
            
            # 获取总数
            count_result = await session.execute(select(func.count(RecordingHistory.id)).where(RecordingHistory.user_id == user_id))
            total = count_result.scalar()
            
            # 分页
            offset = (page - 1) * pageSize
            query = query.offset(offset).limit(pageSize)
            
            result = await session.execute(query)
            recordings = result.scalars().all()
            
            data = []
            for recording in recordings:
                data.append({
                    "id": recording.id,
                    "meeting_id": recording.meeting_id,
                    "meeting_title": recording.meeting_title,
                    "filename": recording.filename,
                    "file_path": recording.file_path,
                    "transcript": recording.transcript[:200] + "..." if len(recording.transcript) > 200 else recording.transcript,
                    "has_minutes": bool(recording.minutes),
                    "duration": recording.duration,
                    "file_size": recording.file_size,
                    "created_at": recording.created_at,
                    "status": recording.status
                })
            
            return {
                "code": 0,
                "data": {
                    "list": data,
                    "total": total,
                    "page": page,
                    "pageSize": pageSize
                }
            }
    
    except Exception as e:
        print(f"获取录音历史失败: {e}")
        return {"code": 1, "msg": f"获取录音历史失败: {str(e)}"}

@app.get("/api/recordings/{recording_id}")
async def get_recording_detail(recording_id: int, payload: dict = Depends(verify_token)):
    """获取录音详情"""
    user_id = payload["userId"]
    
    try:
        async with SessionLocal() as session:
            result = await session.execute(
                select(RecordingHistory).where(
                    RecordingHistory.id == recording_id,
                    RecordingHistory.user_id == user_id
                )
            )
            recording = result.scalar_one_or_none()
            
            if not recording:
                return {"code": 1, "msg": "录音记录不存在"}
            
            return {
                "code": 0,
                "data": {
                    "id": recording.id,
                    "meeting_id": recording.meeting_id,
                    "meeting_title": recording.meeting_title,
                    "filename": recording.filename,
                    "file_path": recording.file_path,
                    "transcript": recording.transcript,
                    "minutes": recording.minutes,
                    "duration": recording.duration,
                    "file_size": recording.file_size,
                    "created_at": recording.created_at,
                    "status": recording.status
                }
            }
    
    except Exception as e:
        print(f"获取录音详情失败: {e}")
        return {"code": 1, "msg": f"获取录音详情失败: {str(e)}"}

@app.put("/api/recordings/{recording_id}")
async def update_recording(recording_id: int, data: RecordingUpdateRequest, payload: dict = Depends(verify_token)):
    """更新录音记录"""
    user_id = payload["userId"]
    
    try:
        async with SessionLocal() as session:
            result = await session.execute(
                select(RecordingHistory).where(
                    RecordingHistory.id == recording_id,
                    RecordingHistory.user_id == user_id
                )
            )
            recording = result.scalar_one_or_none()
            
            if not recording:
                return {"code": 1, "msg": "录音记录不存在"}
            
            # 更新字段
            if data.transcript:
                recording.transcript = data.transcript
            if data.minutes:
                recording.minutes = data.minutes
            
            await session.commit()
            
            return {"code": 0, "msg": "更新成功"}
    
    except Exception as e:
        print(f"更新录音记录失败: {e}")
        return {"code": 1, "msg": f"更新录音记录失败: {str(e)}"}

@app.delete("/api/recordings/{recording_id}")
async def delete_recording(recording_id: int, payload: dict = Depends(verify_token)):
    """删除录音记录"""
    user_id = payload["userId"]
    
    try:
        async with SessionLocal() as session:
            result = await session.execute(
                select(RecordingHistory).where(
                    RecordingHistory.id == recording_id,
                    RecordingHistory.user_id == user_id
                )
            )
            recording = result.scalar_one_or_none()
            
            if not recording:
                return {"code": 1, "msg": "录音记录不存在"}
            
            # 删除文件
            file_full_path = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'uploads', recording.file_path)
            if os.path.exists(file_full_path):
                os.remove(file_full_path)
            
            # 删除数据库记录
            await session.delete(recording)
            await session.commit()
            
            return {"code": 0, "msg": "删除成功"}
    
    except Exception as e:
        print(f"删除录音记录失败: {e}")
        return {"code": 1, "msg": f"删除录音记录失败: {str(e)}"}

@app.get("/api/recordings/{recording_id}/download")
async def download_recording(recording_id: int, payload: dict = Depends(verify_token)):
    """下载录音文件"""
    user_id = payload["userId"]
    
    try:
        async with SessionLocal() as session:
            result = await session.execute(
                select(RecordingHistory).where(
                    RecordingHistory.id == recording_id,
                    RecordingHistory.user_id == user_id
                )
            )
            recording = result.scalar_one_or_none()
            
            if not recording:
                return {"code": 1, "msg": "录音记录不存在"}
            
            file_full_path = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'uploads', recording.file_path)
            
            if not os.path.exists(file_full_path):
                return {"code": 1, "msg": "录音文件不存在"}
            
            return FileResponse(
                file_full_path,
                filename=recording.filename,
                media_type='audio/webm'
            )
    
    except Exception as e:
        print(f"下载录音文件失败: {e}")
        return {"code": 1, "msg": f"下载录音文件失败: {str(e)}"}

# ========== 路由注册 ==========
# 注册各个功能模块的路由
app.include_router(search.router, prefix="/api")
app.include_router(contact.router, prefix="/api")  
app.include_router(group.router, prefix="/api")
app.include_router(message.router, prefix="/api")
app.include_router(assistant_upload.router, prefix="/api/assistant")

# 添加静态文件服务支持
uploads_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), "uploads")
app.mount("/uploads", StaticFiles(directory=uploads_dir), name="uploads")

# ========== 应用启动 ==========
if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=3007)